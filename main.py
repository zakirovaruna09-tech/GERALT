# ══════════════════════════════════════════════════════════════════════
#  GERALT SYSTEM  v5.2  —  Personal AI Assistant
#  Python 3.10+  |  Windows  |  CustomTkinter UI
#  + СТЕГАНОГРАФИЯ: LSB + AES-256 шифрование текста в фото
#  FIX v5.2: исправлена синхронизация TTS → STT (микрофон теперь
#            включается ТОЛЬКО после полного завершения озвучки)
# ══════════════════════════════════════════════════════════════════════

import os, sys, json, time, hashlib, requests, webbrowser
import psutil, subprocess, io, threading, smtplib, random, re, tempfile
import struct
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from datetime import datetime
from colorama import init, Fore, Style
from telebot import TeleBot, types
import customtkinter as ctk
import asyncio, edge_tts

# ══════════════════════════════════════════════════════════════════════
# АУДИО
# ══════════════════════════════════════════════════════════════════════
HAS_VOICE = False
AUDIO_BACKEND = "none"

try:
    from playsound import playsound as _playsound
    HAS_VOICE = True; AUDIO_BACKEND = "playsound"
except ImportError:
    pass

if not HAS_VOICE:
    try:
        import pygame
        pygame.mixer.init()
        HAS_VOICE = True; AUDIO_BACKEND = "pygame"
    except ImportError:
        pass

if not HAS_VOICE:
    try:
        import winsound as _winsound
        HAS_VOICE = True; AUDIO_BACKEND = "winsound"
    except ImportError:
        pass

# МИКРОФОН
HAS_MIC = False; MIC_BACKEND = "none"
try:
    import sounddevice as sd
    import scipy.io.wavfile as _wav_io
    HAS_MIC = True; MIC_BACKEND = "sounddevice"
except ImportError:
    pass

try:
    import speech_recognition as sr
    HAS_SR = True
except ImportError:
    HAS_SR = False

# СТЕГАНОГРАФИЯ
HAS_STEG = False
try:
    from PIL import Image
    from Crypto.Cipher import AES
    from Crypto.Hash import SHA256
    from Crypto.Util.Padding import pad, unpad
    HAS_STEG = True
except ImportError:
    pass

init(autoreset=True)

# ══════════════════════════════════════════════════════════════════════
# КОНФИГУРАЦИЯ — секреты читаются из .env (см. .env.example)
# ══════════════════════════════════════════════════════════════════════
from dotenv import load_dotenv
_BASE_DIR = os.path.dirname(os.path.abspath(__file__))
load_dotenv(os.path.join(_BASE_DIR, ".env"))

BOT_TOKEN          = os.environ.get("BOT_TOKEN_LEGACY", os.environ.get("BOT_TOKEN", ""))
OWNER_ID           = int(os.environ.get("OWNER_ID", "0") or "0")
GMAIL              = os.environ.get("GMAIL", "")
GMAIL_APP_PASSWORD = os.environ.get("GMAIL_APP_PASS_LEGACY", os.environ.get("GMAIL_APP_PASS", ""))
WEATHER_CITY       = "Fergana"
MAX_ATTEMPTS       = 3
LOCK_MINUTES       = 10
PIN_FILE           = os.path.join(_BASE_DIR, "pin.json")
LOG_FILE           = os.path.join(_BASE_DIR, "geralt_log.txt")
MAX_HISTORY        = 10
JARVIS_DIR         = os.environ.get("JARVIS_DIR", "D:\\Jarvis")

if not BOT_TOKEN:
    print("  [CONFIG]  ⚠ BOT_TOKEN не задан. Заполни файл .env (см. .env.example).")
if OWNER_ID == 0:
    print("  [CONFIG]  ⚠ OWNER_ID не задан. Заполни файл .env (см. .env.example).")

# ── FIX v5.2 ──────────────────────────────────────────────────────────
# Пауза ПОСЛЕ завершения воспроизведения аудио перед включением микрофона.
# Увеличено с 0.8 до 1.2 с чтобы избежать самозахвата хвоста озвучки.
TTS_SILENCE_AFTER  = 3.5   # пауза после окончания mp3 (ждём рассеивания звука)
TTS_LISTEN_DELAY   = 1.2   # доп. задержка перед listen() после снятия блокировки
# ─────────────────────────────────────────────────────────────────────

SYSTEM_PROMPT = (
    "Ты Геральт из Ривии — персональный ИИ-ассистент пользователя.\n"
    "КРИТИЧЕСКИ ВАЖНЫЕ ПРАВИЛА:\n"
    "1. ВСЕГДА отвечай ТОЛЬКО на русском языке.\n"
    "2. В КАЖДОМ ответе обращайся ТОЛЬКО 'сэр'. Никаких имён.\n"
    "3. Отвечай кратко, уверенно, в стиле Геральта.\n"
    "4. Помни контекст разговора.\n"
    "Пример: 'Понял вас, сэр. Выполняю.'\n"
)

HF_TOKEN = os.environ.get("HF_TOKEN_LEGACY", os.environ.get("HF_TOKEN", ""))

STEAM_GAMES = {
    "ведьмак":     "steam://rungameid/292030",
    "асасин":      "steam://rungameid/2208920",
    "dead island": "steam://rungameid/1934680",
    "цивилизация": "steam://rungameid/289070",
    "cs2":         "steam://rungameid/730",
}

GERALT_QUOTES = [
    "Зло — это зло. Меньшее, большее, среднее — всё едино.",
    "Если я должен выбирать между одним злом и другим, я предпочитаю не выбирать вовсе.",
    "Мир не нуждается в героях. Он нуждается в профессионалах.",
    "Жизнь коротка. Слишком коротка, чтобы тратить её на сожаления.",
    "Сила без мудрости — лишь разрушение.",
    "Даже в темноте можно найти путь, если знаешь куда идти.",
]

chat_history     = []
bot              = TeleBot(BOT_TOKEN)
user_state       = {}
email_data       = {}
doc_data         = {}
pptx_data        = {}
steg_state       = {}
INTERACTION_MODE = "voice"
app              = None

# ── FIX v5.2: центральный флаг и Event ────────────────────────────────
# IS_SPEAKING = True  →  TTS воспроизводится прямо сейчас
# SPEAK_LOCK  — установлен (set) когда можно слушать, сброшен (clear) когда говорим
IS_SPEAKING   = False
SPEAK_LOCK    = threading.Event()
SPEAK_LOCK.set()   # изначально — можно слушать
# ─────────────────────────────────────────────────────────────────────

_FORBIDDEN_NAMES = ["аруна", "aruna", "пользователь", "друг", "приятель", "господин", "хозяин"]

STEG_DELIMITER = b"<<GERALT_END>>"

# ══════════════════════════════════════════════════════════════════════
# ПАЛИТРА ТЕМ
# ══════════════════════════════════════════════════════════════════════
THEMES = {
    "dark": {
        "bg":           "#0a0b14",
        "surface":      "#10121f",
        "surface2":     "#161928",
        "border":       "#1e2238",
        "border2":      "#2a2f52",
        "accent":       "#7c3aed",
        "accent_hover": "#6d28d9",
        "accent_soft":  "#1e1535",
        "accent2":      "#06b6d4",
        "accent2_soft": "#0c2030",
        "danger":       "#ef4444",
        "danger_soft":  "#2d0f0f",
        "success":      "#22c55e",
        "success_soft": "#0a2010",
        "warn":         "#f59e0b",
        "warn_soft":    "#251a05",
        "text":         "#e2e8f0",
        "text2":        "#94a3b8",
        "text3":        "#475569",
        "text_inv":     "#ffffff",
    },
    "light": {
        "bg":           "#f0f2f8",
        "surface":      "#ffffff",
        "surface2":     "#f8f9fc",
        "border":       "#dde1f0",
        "border2":      "#c8cfe8",
        "accent":       "#6d28d9",
        "accent_hover": "#5b21b6",
        "accent_soft":  "#ede9fe",
        "accent2":      "#0891b2",
        "accent2_soft": "#e0f2fe",
        "danger":       "#dc2626",
        "danger_soft":  "#fee2e2",
        "success":      "#16a34a",
        "success_soft": "#dcfce7",
        "warn":         "#d97706",
        "warn_soft":    "#fef3c7",
        "text":         "#1e293b",
        "text2":        "#475569",
        "text3":        "#94a3b8",
        "text_inv":     "#ffffff",
    }
}

# ══════════════════════════════════════════════════════════════════════
# КОНСОЛЬНЫЕ ХЕЛПЕРЫ
# ══════════════════════════════════════════════════════════════════════
W   = Fore.WHITE
DIM = Fore.WHITE + Style.DIM
M   = Fore.MAGENTA
C   = Fore.CYAN
G   = Fore.GREEN
Y   = Fore.YELLOW
R   = Fore.RED
RS  = Style.RESET_ALL

def _hline(char="═", width=56): return M + char * width + RS
def _sep(width=56):              return M + "╟" + "─" * (width-2) + "╢" + RS

def notify(msg: str, level: str = "info"):
    icons  = {"info":"◈","ok":"✓","warn":"⚠","err":"✗","sys":"⚙"}
    colors = {"info":C,"ok":G,"warn":Y,"err":R,"sys":M}
    icon   = icons.get(level, "·")
    col    = colors.get(level, C)
    now    = datetime.now().strftime("%H:%M:%S")
    print(f"  {col}{icon}{RS} {DIM}{now}{RS}  {W}{msg}{RS}")
    if app:
        app.write_log(f"[{icon}] {msg}")

def print_progress(label: str, value: float, width: int = 20):
    filled  = int(value / 100 * width)
    bar_col = G if value < 60 else (Y if value < 80 else R)
    bar     = bar_col + "█" * filled + DIM + "░" * (width - filled) + RS
    val_col = G if value < 60 else (Y if value < 80 else R)
    print(f"  {DIM}{label:<8}{RS} {bar} {val_col}{value:5.1f}%{RS}")

def _sanitize_response(text: str) -> str:
    result = text
    for name in _FORBIDDEN_NAMES:
        result = re.sub(re.escape(name), "сэр", result, flags=re.IGNORECASE)
    if "сэр" not in result.lower():
        result = "Слушаю, сэр. " + result
    return result

# ══════════════════════════════════════════════════════════════════════
# БЛОК 1 — ЛОГИРОВАНИЕ
# ══════════════════════════════════════════════════════════════════════
def log(text: str, status: str = "OK"):
    line = f"[{datetime.now().strftime('%d.%m.%Y %H:%M:%S')}] [{status:5}] {text}\n"
    try:
        with open(LOG_FILE, "a", encoding="utf-8") as f:
            f.write(line)
    except Exception as e:
        print(f"{R}  ✗ Лог: {e}{RS}")

def show_log(n: int = 15):
    if not os.path.exists(LOG_FILE):
        notify("Лог пуст.", "warn"); return
    with open(LOG_FILE, "r", encoding="utf-8") as f:
        lines = f.readlines()
    print(f"\n{_hline()}")
    for l in lines[-n:]:
        l = l.strip()
        if "BLOCK" in l or "FAIL" in l: print(f"  {R}{l}{RS}")
        elif "WARN" in l:               print(f"  {Y}{l}{RS}")
        else:                           print(f"  {DIM}{l}{RS}")
    print(_hline() + "\n")

# ══════════════════════════════════════════════════════════════════════
# БЛОК 2 — PIN
# ══════════════════════════════════════════════════════════════════════
def _hash(pin: str) -> str:
    return hashlib.sha256(pin.encode()).hexdigest()

def _pin_exists() -> bool: return os.path.exists(PIN_FILE)

def _save_pin(pin: str):
    with open(PIN_FILE, "w") as f:
        json.dump({"hash": _hash(pin)}, f)

def _load_hash() -> str:
    with open(PIN_FILE, "r") as f:
        return json.load(f)["hash"]

def _setup_pin():
    print(f"\n{_hline()}")
    print(f"  {Y}ПЕРВЫЙ ЗАПУСК — УСТАНОВКА PIN{RS}")
    print(_sep())
    while True:
        pin = input(f"  {C}Новый PIN (4 цифры): {W}").strip()
        if not pin.isdigit() or len(pin) != 4:
            notify("PIN — ровно 4 цифры.", "err"); continue
        confirm = input(f"  {C}Повторите PIN: {W}").strip()
        if pin != confirm:
            notify("Не совпадает.", "err"); continue
        _save_pin(pin)
        notify("PIN установлен!", "ok")
        log("PIN установлен", "INFO"); break
    print(_hline() + "\n")

def change_pin():
    print(f"\n{_hline()}")
    print(f"  {C}СМЕНА PIN-КОДА{RS}")
    print(_sep())
    old = input(f"  {Y}Текущий PIN: {W}").strip()
    if _hash(old) != _load_hash():
        notify("Неверный PIN.", "err")
        log("Смена PIN — неверный текущий", "WARN")
        print(_hline() + "\n"); return False
    while True:
        new_pin = input(f"  {C}Новый PIN (4 цифры): {W}").strip()
        if not new_pin.isdigit() or len(new_pin) != 4:
            notify("PIN — ровно 4 цифры.", "err"); continue
        confirm = input(f"  {C}Повторите: {W}").strip()
        if new_pin != confirm:
            notify("Не совпадает.", "err"); continue
        _save_pin(new_pin)
        notify("PIN изменён!", "ok")
        log("PIN изменён", "INFO")
        print(_hline() + "\n"); return True

# ══════════════════════════════════════════════════════════════════════
# БЛОК 3 — АВТОБЛОКИРОВКА
# ══════════════════════════════════════════════════════════════════════
def _lock():
    lock_until = time.time() + LOCK_MINUTES * 60
    log(f"ЗАБЛОКИРОВАНО на {LOCK_MINUTES} мин.", "BLOCK")
    print(f"\n{_hline()}")
    print(f"  {R}СИСТЕМА ЗАБЛОКИРОВАНА НА {LOCK_MINUTES} МИНУТ{RS}")
    print(_hline())
    while time.time() < lock_until:
        remaining = int(lock_until - time.time())
        m, s = remaining // 60, remaining % 60
        print(f"\r  {R}Разблокировка через: {Y}{m:02d}:{s:02d}{RS}   ", end="", flush=True)
        time.sleep(1)
    print(f"\n\n  {G}Блокировка снята.{RS}\n")
    log("Блокировка снята", "INFO")

def security_check() -> bool:
    if not _pin_exists():
        _setup_pin()
    attempts = 0
    print(f"\n{_hline()}")
    print(f"  {M}GERALT SYSTEM  ·  {C}ПРОВЕРКА PIN{RS}")
    print(_sep())
    while attempts < MAX_ATTEMPTS:
        left = MAX_ATTEMPTS - attempts
        dots = f"{G}●{RS}" * left + f"{DIM}●{RS}" * (MAX_ATTEMPTS - left)
        pin = input(f"  {Y}PIN [{dots}{Y}]: {W}").strip()
        if not pin.isdigit() or len(pin) != 4:
            notify("PIN — ровно 4 цифры.", "warn"); continue
        if _hash(pin) == _load_hash():
            notify("Доступ разрешён. Добро пожаловать, сэр.", "ok")
            print(_hline() + "\n")
            log("Вход по PIN", "OK"); return True
        attempts += 1
        remaining = MAX_ATTEMPTS - attempts
        log(f"Неверный PIN ({attempts}/{MAX_ATTEMPTS})", "FAIL")
        if remaining > 0:
            notify(f"Неверный PIN. Осталось: {remaining}", "err")
        else:
            notify("Неверный PIN. Блокировка.", "err")
    print(_hline() + "\n")
    _lock()
    return security_check()

# ══════════════════════════════════════════════════════════════════════
# БЛОК 4 — TELEGRAM
# ══════════════════════════════════════════════════════════════════════
def only_owner(func):
    def wrapper(message):
        if message.chat.id != OWNER_ID:
            bot.send_message(message.chat.id, "⛔ Доступ запрещён.")
            log(f"Чужой доступ: {message.chat.id}", "WARN"); return
        return func(message)
    return wrapper

@bot.message_handler(commands=["start"])
@only_owner
def tg_start(message):
    markup = types.ReplyKeyboardMarkup(resize_keyboard=True)
    markup.row("💻 Статус ПК", "🔌 Выключить ПК")
    markup.row("📧 Отправить Email", "💱 Курс валют")
    markup.row("☀️ Сводка дня", "📄 Текст в документ")
    markup.row("📊 Презентация", "❓ Помощь")
    markup.row("🔒 Скрыть текст", "🔓 Извлечь текст")
    bot.send_message(message.chat.id,
        "👋 *GERALT SYSTEM v5.2*\n\nВыбери команду:",
        parse_mode="Markdown", reply_markup=markup)

@bot.message_handler(func=lambda m: m.text == "💻 Статус ПК")
@only_owner
def pc_status(message):
    cpu  = psutil.cpu_percent(interval=1)
    ram  = psutil.virtual_memory()
    disk = psutil.disk_usage("C:\\")
    uptime_sec = int(time.time() - psutil.boot_time())
    h, mn = uptime_sec // 3600, (uptime_sec % 3600) // 60
    text = (
        f"🖥 *Статус компьютера*\n\n"
        f"⚙️ CPU: `{cpu}%`\n"
        f"🧠 RAM: `{ram.percent}%` ({round(ram.used/1024**3,1)} / {round(ram.total/1024**3,1)} ГБ)\n"
        f"💾 Диск C: `{disk.percent}%` ({round(disk.used/1024**3,1)} / {round(disk.total/1024**3,1)} ГБ)\n"
        f"⏱ Аптайм: `{h}ч {mn}м`\n✅ Онлайн"
    )
    bot.send_message(message.chat.id, text, parse_mode="Markdown")

@bot.message_handler(func=lambda m: m.text == "🔌 Выключить ПК")
@only_owner
def shutdown_confirm(message):
    markup = types.InlineKeyboardMarkup()
    markup.row(
        types.InlineKeyboardButton("✅ Да, выключить", callback_data="shutdown_yes"),
        types.InlineKeyboardButton("❌ Отмена",         callback_data="shutdown_no")
    )
    bot.send_message(message.chat.id, "⚠️ *Ты уверен?*",
                     parse_mode="Markdown", reply_markup=markup)

@bot.callback_query_handler(func=lambda c: c.data in ["shutdown_yes", "shutdown_no"])
def shutdown_callback(call):
    if call.from_user.id != OWNER_ID: return
    if call.data == "shutdown_yes":
        bot.answer_callback_query(call.id, "Выключаю...")
        bot.edit_message_text("🔌 *Выключение через 5 сек...*",
                              call.message.chat.id, call.message.message_id, parse_mode="Markdown")
        log("Выключение через TG", "WARN")
        os.system("shutdown /s /t 5")
    else:
        bot.answer_callback_query(call.id, "Отменено")
        bot.edit_message_text("❌ *Отменено.*",
                              call.message.chat.id, call.message.message_id, parse_mode="Markdown")

@bot.message_handler(func=lambda m: m.text == "📧 Отправить Email")
@only_owner
def email_start(message):
    user_state[message.chat.id] = "waiting_to"
    email_data[message.chat.id] = {}
    bot.send_message(message.chat.id, "📧 Введи адрес получателя:")

@bot.message_handler(func=lambda m: user_state.get(m.chat.id) == "waiting_to")
@only_owner
def email_get_to(message):
    email_data[message.chat.id]["to"] = message.text
    user_state[message.chat.id] = "waiting_subject"
    bot.send_message(message.chat.id, "📝 Тема письма:")

@bot.message_handler(func=lambda m: user_state.get(m.chat.id) == "waiting_subject")
@only_owner
def email_get_subject(message):
    email_data[message.chat.id]["subject"] = message.text
    user_state[message.chat.id] = "waiting_body"
    bot.send_message(message.chat.id, "✏️ Текст письма:")

@bot.message_handler(func=lambda m: user_state.get(m.chat.id) == "waiting_body")
@only_owner
def email_get_body(message):
    email_data[message.chat.id]["body"] = message.text
    user_state[message.chat.id] = None
    data = email_data[message.chat.id]
    bot.send_message(message.chat.id, "📤 Отправляю...")
    try:
        msg = MIMEMultipart()
        msg["From"]    = GMAIL
        msg["To"]      = data["to"]
        msg["Subject"] = data["subject"]
        msg.attach(MIMEText(data["body"], "plain"))
        with smtplib.SMTP_SSL("smtp.gmail.com", 465) as server:
            server.login(GMAIL, GMAIL_APP_PASSWORD)
            server.sendmail(GMAIL, data["to"], msg.as_string())
        bot.send_message(message.chat.id,
            f"✅ *Отправлено!*\n📬 Кому: `{data['to']}`\n📝 Тема: `{data['subject']}`",
            parse_mode="Markdown")
        log(f"Email → {data['to']}", "OK")
    except Exception as e:
        bot.send_message(message.chat.id, f"❌ Ошибка: {e}")
        log(f"Ошибка email: {e}", "WARN")

@bot.message_handler(func=lambda m: m.text == "💱 Курс валют")
@only_owner
def currency_cmd(message):
    rates = get_currency_rates()
    if rates:
        text = (f"💱 *Курс валют*\n\n"
                f"🇺🇸 USD → UZS: `{rates['usd_uzs']:,.0f} сум`\n"
                f"🇷🇺 RUB → UZS: `{rates['rub_uzs']:,.2f} сум`\n"
                f"💵 USD → RUB: `{rates['usd_rub']:,.2f} руб`\n\n"
                f"🕐 {datetime.now().strftime('%H:%M')}")
    else:
        text = "❌ Нет данных."
    bot.send_message(message.chat.id, text, parse_mode="Markdown")

@bot.message_handler(func=lambda m: m.text == "☀️ Сводка дня")
@only_owner
def briefing_cmd(message):
    bot.send_message(message.chat.id, build_morning_briefing(), parse_mode="Markdown")

@bot.message_handler(func=lambda m: m.text == "📄 Текст в документ")
@only_owner
def doc_start(message):
    user_state[message.chat.id] = "doc_waiting_title"
    doc_data[message.chat.id]   = {}
    bot.send_message(message.chat.id, "📄 *Создание .docx*\n\nНазвание документа:", parse_mode="Markdown")

@bot.message_handler(func=lambda m: user_state.get(m.chat.id) == "doc_waiting_title")
@only_owner
def doc_get_title(message):
    doc_data[message.chat.id]["title"] = message.text
    user_state[message.chat.id] = "doc_waiting_text"
    bot.send_message(message.chat.id, "✏️ Текст документа:")

@bot.message_handler(func=lambda m: user_state.get(m.chat.id) == "doc_waiting_text")
@only_owner
def doc_get_text(message):
    chat_id = message.chat.id
    doc_data[chat_id]["text"] = message.text
    user_state[chat_id] = None
    title = doc_data[chat_id]["title"]
    text  = doc_data[chat_id]["text"]
    bot.send_message(chat_id, "⏳ Создаю...")
    filepath = _create_docx(title, text)
    if filepath:
        with open(filepath, "rb") as f:
            bot.send_document(chat_id, f,
                caption=f"✅ *Готово, сэр!*\n📋 `{title}`", parse_mode="Markdown")
        os.remove(filepath)
        log(f"Документ TG: {title}", "OK")
    else:
        bot.send_message(chat_id, "❌ Ошибка. pip install python-docx")
    doc_data.pop(chat_id, None)

@bot.message_handler(func=lambda m: m.text == "📊 Презентация")
@only_owner
def pptx_start(message):
    user_state[message.chat.id] = "pptx_waiting_topic"
    pptx_data[message.chat.id]  = {}
    bot.send_message(message.chat.id, "📊 *Создание .pptx*\n\nТема:", parse_mode="Markdown")

@bot.message_handler(func=lambda m: user_state.get(m.chat.id) == "pptx_waiting_topic")
@only_owner
def pptx_get_topic(message):
    pptx_data[message.chat.id]["topic"] = message.text
    user_state[message.chat.id] = "pptx_waiting_slides"
    bot.send_message(message.chat.id, "🔢 Сколько слайдов? (3–10):")

@bot.message_handler(func=lambda m: user_state.get(m.chat.id) == "pptx_waiting_slides")
@only_owner
def pptx_get_slides(message):
    chat_id = message.chat.id
    try:
        count = int(message.text.strip())
        if not (3 <= count <= 10): raise ValueError
    except:
        bot.send_message(chat_id, "❌ Введи число от 3 до 10:"); return
    pptx_data[chat_id]["slides"] = count
    user_state[chat_id] = None
    topic = pptx_data[chat_id]["topic"]
    bot.send_message(chat_id, f"⏳ Генерирую: *{topic}*...", parse_mode="Markdown")
    filepath = generate_pptx(topic, count)
    if filepath and os.path.exists(filepath):
        with open(filepath, "rb") as f:
            bot.send_document(chat_id, f,
                caption=f"✅ *Готово!*\n📋 `{topic}` · {count} слайдов", parse_mode="Markdown")
        os.remove(filepath)
        log(f"Презентация TG: {topic}", "OK")
    else:
        bot.send_message(chat_id, "❌ Ошибка. pip install python-pptx")
    pptx_data.pop(chat_id, None)

@bot.message_handler(func=lambda m: m.text == "❓ Помощь")
@only_owner
def help_cmd(message):
    steg_status = "✅ активна" if HAS_STEG else "❌ pip install Pillow pycryptodome"
    bot.send_message(message.chat.id,
        "📋 *Команды:*\n\n"
        "💻 Статус ПК · 🔌 Выключить ПК\n"
        "📧 Email · 💱 Курс валют\n"
        "☀️ Сводка дня · 📄 Документ · 📊 Презентация\n\n"
        f"🔒 Скрыть текст · 🔓 Извлечь текст\n"
        f"_(стеганография: {steg_status})_",
        parse_mode="Markdown")

# ══════════════════════════════════════════════════════════════════════
# БЛОК 4B — TELEGRAM СТЕГАНОГРАФИЯ
# ══════════════════════════════════════════════════════════════════════
@bot.message_handler(func=lambda m: m.text in ("🔒 Скрыть текст", "🔓 Извлечь текст"))
@only_owner
def steg_menu(message):
    chat_id = message.chat.id
    if not HAS_STEG:
        bot.send_message(chat_id,
            "❌ Библиотеки не установлены.\n\n"
            "Установи:\n`pip install Pillow pycryptodome`",
            parse_mode="Markdown")
        return
    if message.text == "🔒 Скрыть текст":
        steg_state[chat_id] = {"mode": "embed", "step": "waiting_photo"}
        bot.send_message(chat_id,
            "🔒 *Скрытие текста в фото*\n\n"
            "Шаг 1/3 — Отправь фотографию.\n"
            "⚠️ Лучше отправляй как *документ* (скрепка), чтобы не потерять качество.",
            parse_mode="Markdown")
    else:
        steg_state[chat_id] = {"mode": "extract", "step": "waiting_photo"}
        bot.send_message(chat_id,
            "🔓 *Извлечение текста из фото*\n\n"
            "Шаг 1/2 — Отправь фото со скрытым текстом как *документ*.",
            parse_mode="Markdown")

@bot.message_handler(content_types=["photo"],
                     func=lambda m: steg_state.get(m.chat.id, {}).get("step") in ("waiting_photo",))
@only_owner
def steg_photo_as_photo(message):
    chat_id = message.chat.id
    state = steg_state.get(chat_id, {})
    if not state: return
    file_info = bot.get_file(message.photo[-1].file_id)
    downloaded = bot.download_file(file_info.file_path)
    tmp_in = os.path.join(tempfile.gettempdir(), f"steg_in_{chat_id}.jpg")
    with open(tmp_in, "wb") as f:
        f.write(downloaded)
    state["photo_path"] = tmp_in
    _steg_after_photo(chat_id, state)

@bot.message_handler(content_types=["document"],
                     func=lambda m: steg_state.get(m.chat.id, {}).get("step") in ("waiting_photo",))
@only_owner
def steg_photo_as_doc(message):
    chat_id = message.chat.id
    state = steg_state.get(chat_id, {})
    if not state: return
    doc = message.document
    mime = doc.mime_type or ""
    if not mime.startswith("image/"):
        bot.send_message(chat_id, "❌ Отправь файл изображения (PNG/JPG)."); return
    file_info = bot.get_file(doc.file_id)
    downloaded = bot.download_file(file_info.file_path)
    ext = ".png" if "png" in mime else ".jpg"
    tmp_in = os.path.join(tempfile.gettempdir(), f"steg_in_{chat_id}{ext}")
    with open(tmp_in, "wb") as f:
        f.write(downloaded)
    state["photo_path"] = tmp_in
    _steg_after_photo(chat_id, state)

def _steg_after_photo(chat_id: int, state: dict):
    if state["mode"] == "embed":
        state["step"] = "waiting_text"
        bot.send_message(chat_id,
            "✅ Фото получено.\n\nШаг 2/3 — Введи *текст* для скрытия:",
            parse_mode="Markdown")
    else:
        state["step"] = "waiting_password_extract"
        bot.send_message(chat_id,
            "✅ Фото получено.\n\nШаг 2/2 — Введи *пароль* для расшифровки:",
            parse_mode="Markdown")

@bot.message_handler(func=lambda m: steg_state.get(m.chat.id, {}).get("step") == "waiting_text")
@only_owner
def steg_get_text(message):
    chat_id = message.chat.id
    steg_state[chat_id]["secret_text"] = message.text
    steg_state[chat_id]["step"] = "waiting_password_embed"
    bot.send_message(chat_id,
        "✅ Текст получен.\n\nШаг 3/3 — Введи *пароль*\n"
        "_(запомни его — без него расшифровать невозможно!)_:",
        parse_mode="Markdown")

@bot.message_handler(func=lambda m: steg_state.get(m.chat.id, {}).get("step") == "waiting_password_embed")
@only_owner
def steg_do_embed(message):
    chat_id = message.chat.id
    state = steg_state.pop(chat_id, {})
    password = message.text.strip()
    if not password:
        bot.send_message(chat_id, "❌ Пароль не может быть пустым."); return

    bot.send_message(chat_id, "⏳ Шифрую и встраиваю в пиксели...")

    def _run():
        try:
            encrypted = steg_encrypt_text(state["secret_text"], password)
            os.makedirs(JARVIS_DIR, exist_ok=True)
            out_path = os.path.join(
                JARVIS_DIR,
                f"steg_out_{datetime.now().strftime('%d%m%Y_%H%M%S')}.png"
            )
            success = steg_embed(state["photo_path"], encrypted, out_path)
            if success and os.path.exists(out_path):
                with open(out_path, "rb") as f:
                    bot.send_document(chat_id, f,
                        caption=(
                            "✅ *Готово, сэр!*\n\n"
                            "Текст зашифрован и скрыт внутри фото.\n\n"
                            "⚠️ *Важно:* пересылай этот файл ТОЛЬКО как документ,\n"
                            "иначе Telegram пересожмёт фото и данные потеряются.\n\n"
                            f"🔑 Пароль: `{password}`"
                        ),
                        parse_mode="Markdown")
                os.remove(out_path)
                log("Стеганография: текст скрыт", "OK")
            else:
                bot.send_message(chat_id,
                    "❌ Ошибка: фото слишком маленькое.\n"
                    "Используй фото побольше или сократи текст.")
        except Exception as e:
            bot.send_message(chat_id, f"❌ Ошибка: {e}")
            log(f"steg embed: {e}", "WARN")
        finally:
            try: os.remove(state.get("photo_path", ""))
            except: pass

    threading.Thread(target=_run, daemon=True).start()

@bot.message_handler(func=lambda m: steg_state.get(m.chat.id, {}).get("step") == "waiting_password_extract")
@only_owner
def steg_do_extract(message):
    chat_id = message.chat.id
    state = steg_state.pop(chat_id, {})
    password = message.text.strip()

    bot.send_message(chat_id, "⏳ Извлекаю и расшифровываю...")

    def _run():
        try:
            result = steg_extract(state["photo_path"], password)
            if result is not None:
                if len(result) <= 3900:
                    bot.send_message(chat_id,
                        f"✅ *Скрытое сообщение:*\n\n`{result}`",
                        parse_mode="Markdown")
                else:
                    bot.send_message(chat_id, "✅ *Скрытое сообщение:*", parse_mode="Markdown")
                    for i in range(0, len(result), 3900):
                        bot.send_message(chat_id, result[i:i+3900])
                log("Стеганография: текст извлечён", "OK")
            else:
                bot.send_message(chat_id,
                    "❌ *Не удалось извлечь текст.*\n\n"
                    "Возможные причины:\n"
                    "• Неверный пароль\n"
                    "• Фото не содержит скрытого текста\n"
                    "• Фото было пересжато (JPEG артефакты от Telegram)\n\n"
                    "_Совет: отправляй фото как документ (скрепка), не как фото._",
                    parse_mode="Markdown")
        except Exception as e:
            bot.send_message(chat_id, f"❌ Ошибка: {e}")
            log(f"steg extract: {e}", "WARN")
        finally:
            try: os.remove(state.get("photo_path", ""))
            except: pass

    threading.Thread(target=_run, daemon=True).start()

# ══════════════════════════════════════════════════════════════════════
# БЛОК 5 — ДАННЫЕ
# ══════════════════════════════════════════════════════════════════════
def get_currency_rates():
    try:
        resp  = requests.get("https://api.exchangerate-api.com/v4/latest/USD", timeout=5)
        rates = resp.json().get("rates", {})
        uzs, rub = rates.get("UZS"), rates.get("RUB")
        if not uzs or not rub: return None
        return {"usd_uzs": round(uzs,0), "rub_uzs": round(uzs/rub,2), "usd_rub": round(rub,2)}
    except Exception as e:
        notify(f"Курс валют: {e}", "err"); return None

def get_weather():
    try:
        resp = requests.get(f"https://wttr.in/{WEATHER_CITY}?format=j1", timeout=10)
        if resp.status_code == 200:
            cond  = resp.json().get("current_condition", [{}])[0]
            temp  = cond.get("temp_C", "?")
            feels = cond.get("FeelsLikeC", "?")
            hum   = cond.get("humidity", "?")
            desc_list = cond.get("weatherDesc", [])
            desc  = desc_list[0]["value"] if desc_list else "Ясно"
            icons = {"sunny":"☀","clear":"☾","cloudy":"☁","overcast":"☁","rain":"⛆",
                     "snow":"❄","fog":"≋","thunder":"⛈","mist":"≋","partly":"⛅"}
            icon  = next((v for k,v in icons.items() if k in desc.lower()), "⊙")
            return {"temp":temp, "feels":feels, "desc":desc, "humidity":hum, "icon":icon}
    except: pass
    try:
        url  = "https://api.open-meteo.com/v1/forecast?latitude=40.38&longitude=71.78&current_weather=true&hourly=relativehumidity_2m"
        data = requests.get(url, timeout=10).json()
        cw   = data.get("current_weather", {})
        wc   = cw.get("weathercode", 0)
        dmap = {0:"Ясно",1:"Преимущественно ясно",2:"Переменная облачность",
                3:"Пасмурно",45:"Туман",51:"Морось",61:"Дождь",71:"Снег",80:"Ливень",95:"Гроза"}
        return {"temp": str(int(cw.get("temperature",0))), "feels": str(int(cw.get("temperature",0))),
                "desc": dmap.get(wc,"Переменная облачность"),
                "humidity": str(data.get("hourly",{}).get("relativehumidity_2m",[50])[0]),
                "icon": "🌡"}
    except Exception as e:
        notify(f"Погода: {e}", "err"); return None

def build_morning_briefing():
    now      = datetime.now()
    greeting = "Доброе утро" if now.hour < 12 else ("Добрый день" if now.hour < 17 else "Добрый вечер")
    text  = f"⚔️ *{greeting}, сэр!*\n📅 {now.strftime('%d.%m.%Y')} | {now.strftime('%H:%M')}\n\n"
    w = get_weather()
    if w:
        text += (f"{w['icon']} *Погода — Фергана*\n"
                 f"🌡 {w['temp']}°C (ощущается {w['feels']}°C)\n"
                 f"☁️ {w['desc']} · 💧 {w['humidity']}%\n\n")
    r = get_currency_rates()
    if r:
        text += (f"💱 *Курс валют*\n"
                 f"🇺🇸 USD: `{r['usd_uzs']:,.0f}` сум\n"
                 f"🇷🇺 RUB: `{r['rub_uzs']:,.2f}` сум\n\n")
    text += f"⚔️ *Геральт:*\n_{random.choice(GERALT_QUOTES)}_"
    return text

# ══════════════════════════════════════════════════════════════════════
# БЛОК 6 — БЕЗОПАСНОСТЬ / НАРУШИТЕЛЬ
# ══════════════════════════════════════════════════════════════════════
def notify_intruder(frame_bgr):
    try:
        import cv2
        timestamp = datetime.now().strftime("%d.%m.%Y %H:%M:%S")
        success, buffer = cv2.imencode(".jpg", frame_bgr)
        if not success: return
        photo_bytes = io.BytesIO(buffer.tobytes())
        photo_bytes.name = "intruder.jpg"
        bot.send_photo(OWNER_ID, photo_bytes,
            caption=f"🚨 *ВЗЛОМ FACE ID*\n🕐 `{timestamp}`\n❌ Отклонён",
            parse_mode="Markdown")
        log(f"Нарушитель — {timestamp}", "WARN")
    except Exception as e:
        notify(f"Безопасность: {e}", "err")

# ══════════════════════════════════════════════════════════════════════
# БЛОК 7 — СОЗДАНИЕ ДОКУМЕНТА
# ══════════════════════════════════════════════════════════════════════
def _create_docx(title: str, text: str):
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        doc = Document()
        heading = doc.add_heading(title, level=1)
        heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in heading.runs:
            run.font.color.rgb = RGBColor(0x2E, 0x74, 0xB5)
        dp = doc.add_paragraph()
        dp.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        dr = dp.add_run(f"Дата: {datetime.now().strftime('%d.%m.%Y')}")
        dr.font.size = Pt(10)
        dr.font.color.rgb = RGBColor(0x70, 0x70, 0x70)
        doc.add_paragraph()
        for para in text.split("\n"):
            if para.strip():
                p = doc.add_paragraph(para.strip())
                p.paragraph_format.space_after = Pt(6)
        os.makedirs(JARVIS_DIR, exist_ok=True)
        filename = f"{title}_{datetime.now().strftime('%d%m%Y_%H%M%S')}.docx"
        filepath = os.path.join(JARVIS_DIR, filename)
        doc.save(filepath)
        return filepath
    except ImportError:
        notify("pip install python-docx", "err"); return None
    except Exception as e:
        notify(f"Ошибка docx: {e}", "err"); return None

# ══════════════════════════════════════════════════════════════════════
# БЛОК 8 — ПРЕЗЕНТАЦИЯ
# ══════════════════════════════════════════════════════════════════════
def generate_pptx(topic: str, n_slides: int = 5):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor as PptxRGB
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        notify("pip install python-pptx", "err"); return None

    notify("Запрашиваю структуру...", "sys")
    slides_data = []
    try:
        import g4f
        prompt = (
            f"Создай структуру презентации на тему: '{topic}'. "
            f"Ровно {n_slides} слайдов. Тексты ТОЛЬКО на русском. "
            f"Отвечай ТОЛЬКО JSON без markdown:\n"
            f'[{{"title":"...","points":["...","...","..."]}}]'
        )
        raw = str(g4f.ChatCompletion.create(
            model=g4f.models.gpt_4,
            messages=[{"role":"user","content":prompt}],
            stream=False
        )).strip()
        if "```" in raw:
            raw = raw.split("```")[1]
            if raw.startswith("json"): raw = raw[4:]
        import json as _json
        slides_data = _json.loads(raw.strip())
    except Exception as e:
        notify(f"GPT заглушка: {e}", "warn")
        slides_data = [{"title": f"Слайд {i+1}", "points": [f"Пункт {j+1}" for j in range(3)]} for i in range(n_slides)]

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    BG    = PptxRGB(0x1A, 0x1A, 0x2E)
    ACCT  = PptxRGB(0xE9, 0x4F, 0x37)
    WHITE = PptxRGB(0xFF, 0xFF, 0xFF)
    GRAY  = PptxRGB(0xCC, 0xCC, 0xCC)

    def set_bg(slide):
        fill = slide.background.fill
        fill.solid(); fill.fore_color.rgb = BG

    def add_tb(slide, text, l, t, w, h, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True
        p  = tf.paragraphs[0]; p.alignment = align
        run = p.add_run(); run.text = text
        run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color

    def add_bar(slide, l, t, w, h, color):
        bar = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s); add_bar(s, 0, 3.2, 13.33, 0.08, ACCT)
    add_tb(s, topic.upper(), 1, 1.5, 11, 1.5, 44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_tb(s, f"Презентация · {datetime.now().strftime('%d.%m.%Y')}", 1, 3.5, 11, 0.6, 18,
           color=PptxRGB(0x99,0x99,0x99), align=PP_ALIGN.CENTER)

    for i, sd_item in enumerate(slides_data[:n_slides]):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(s); add_bar(s, 0, 0, 13.33, 0.08, ACCT)
        add_tb(s, f"{i+1:02d}", 11.8, 0.15, 1.2, 0.5, 14, color=PptxRGB(0x66,0x66,0x66))
        add_tb(s, sd_item.get("title", f"Слайд {i+1}"), 0.5, 0.3, 11, 1.0, 32, bold=True, color=WHITE)
        add_bar(s, 0.5, 1.4, 12.3, 0.03, PptxRGB(0x44,0x44,0x66))
        y = 1.6
        for point in sd_item.get("points", [])[:6]:
            add_tb(s, f"▸  {point}", 0.7, y, 11.5, 0.7, 20, color=GRAY); y += 0.75

    os.makedirs(JARVIS_DIR, exist_ok=True)
    filename = f"geralt_pptx_{datetime.now().strftime('%d%m%Y_%H%M%S')}.pptx"
    filepath = os.path.join(JARVIS_DIR, filename)
    prs.save(filepath)
    notify(f"Презентация → {filename}", "ok")
    log(f"Презентация: {filename}", "OK")
    return filepath

# ══════════════════════════════════════════════════════════════════════
# БЛОК 9 — ПЛАНИРОВЩИК
# ══════════════════════════════════════════════════════════════════════
def scheduler_loop():
    morning_sent_date  = None
    pc_start_time      = datetime.now()
    last_rest_interval = 0
    notify("Планировщик запущен.", "sys")
    while True:
        now = datetime.now()
        if now.hour == 8 and now.minute == 0:
            today = now.date()
            if morning_sent_date != today:
                morning_sent_date = today
                try:
                    bot.send_message(OWNER_ID, build_morning_briefing(), parse_mode="Markdown")
                    log("Утренняя сводка отправлена", "OK")
                except Exception as e:
                    notify(f"Планировщик: {e}", "err")
        hours_on = (datetime.now() - pc_start_time).total_seconds() / 3600
        current_interval = int(hours_on // 4)
        if current_interval > 0 and current_interval != last_rest_interval:
            last_rest_interval = current_interval
            try:
                bot.send_message(OWNER_ID,
                    f"⏰ *Сэр, вы за компом {int(hours_on)} ч.*\nСоветую перерыв.\n_— Геральт_",
                    parse_mode="Markdown")
            except: pass
        time.sleep(60)

# ══════════════════════════════════════════════════════════════════════
# БЛОК 10 — ГОЛОС TTS  (FIX v5.2)
# ══════════════════════════════════════════════════════════════════════
async def _tts_generate(text: str, filename: str):
    communicate = edge_tts.Communicate(text, "ru-RU-DmitryNeural")
    await communicate.save(filename)

def _play_audio(filepath: str):
    """Воспроизводит аудио-файл. Возвращает управление только после окончания."""
    global AUDIO_BACKEND, HAS_VOICE
    if AUDIO_BACKEND == "playsound":
        try:
            _playsound(filepath, block=True)   # block=True — ждём конца
            return
        except Exception as e:
            notify(f"playsound: {e}", "warn"); AUDIO_BACKEND = "fallback"

    if AUDIO_BACKEND in ("pygame", "pygame_ce"):
        try:
            import pygame
            if not pygame.mixer.get_init(): pygame.mixer.init()
            pygame.mixer.music.load(filepath)
            pygame.mixer.music.play()
            # FIX: ждём явного окончания воспроизведения
            while pygame.mixer.music.get_busy():
                time.sleep(0.05)
            return
        except Exception as e:
            notify(f"pygame: {e}", "warn"); AUDIO_BACKEND = "fallback"

    if AUDIO_BACKEND in ("winsound", "fallback"):
        try:
            wav_path = filepath.replace(".mp3", ".wav")
            result = subprocess.run(["ffmpeg", "-y", "-i", filepath, wav_path],
                                    capture_output=True, timeout=15)
            if result.returncode == 0 and os.path.exists(wav_path):
                import winsound
                winsound.PlaySound(wav_path, winsound.SND_FILENAME)
                os.remove(wav_path)
                return
        except FileNotFoundError:
            notify("ffmpeg не найден → os.startfile", "warn")
        except Exception as e:
            notify(f"winsound: {e}", "warn")

    # Последний резерв
    try:
        os.startfile(filepath)
        time.sleep(4)   # минимальная пауза, т.к. startfile не блокирует
    except Exception as e:
        notify(f"os.startfile: {e}", "err")


def _tts_thread(text: str):
    """
    FIX v5.2 — правильная последовательность:
    1. IS_SPEAKING = True  + SPEAK_LOCK сброшен  (микрофон заглушён)
    2. Генерация mp3
    3. Воспроизведение (block=True / get_busy-loop)
    4. TTS_SILENCE_AFTER — доп. тишина после конца аудио
    5. IS_SPEAKING = False + SPEAK_LOCK установлен  (микрофон разрешён)
    """
    global IS_SPEAKING

    tmp = tempfile.NamedTemporaryFile(suffix=".mp3", delete=False, prefix="geralt_tts_")
    tmp_path = tmp.name
    tmp.close()

    try:
        # ── 1. Сигнализируем: идёт озвучка ─────────────────────────
        IS_SPEAKING = True
        SPEAK_LOCK.clear()          # блокируем listen() в voice_loop

        # ── 2. Генерация TTS ────────────────────────────────────────
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        loop.run_until_complete(_tts_generate(text, tmp_path))
        loop.close()

        # ── 3. Воспроизведение (блокирующее) ────────────────────────
        _play_audio(tmp_path)

        # ── 4. Пауза после конца аудио (звук рассеивается) ──────────
        time.sleep(TTS_SILENCE_AFTER)

    except Exception as e:
        notify(f"TTS: {e}", "err")

    finally:
        # ── 5. Разрешаем микрофон ТОЛЬКО ЗДЕСЬ ─────────────────────
        IS_SPEAKING = False
        SPEAK_LOCK.set()

        try:
            if os.path.exists(tmp_path):
                os.remove(tmp_path)
        except:
            pass


def speak(text: str):
    now = datetime.now().strftime("%H:%M:%S")
    print(f"\n  {DIM}{now}{RS}  {M}⚔{RS}  {W}{text}{RS}\n")
    if app:
        app.write_log(f"Геральт: {text}", tag="geralt")
    if HAS_VOICE and INTERACTION_MODE == "voice":
        threading.Thread(target=_tts_thread, args=(text,), daemon=True).start()

# ══════════════════════════════════════════════════════════════════════
# БЛОК 11 — МИКРОФОН STT
# ══════════════════════════════════════════════════════════════════════
def listen() -> str:
    if not HAS_SR:
        notify("speech_recognition не установлен", "err"); return ""
    recognizer = sr.Recognizer()
    if HAS_MIC and MIC_BACKEND == "sounddevice":
        try:
            print(f"\n  {M}◉ {C}Слушаю... (4 сек){RS}")
            fs = 16000; duration = 4
            recording = sd.rec(int(duration * fs), samplerate=fs, channels=1,
                               dtype='int16', blocking=True)
            byte_io = io.BytesIO()
            _wav_io.write(byte_io, fs, recording)
            byte_io.seek(0)
            with sr.AudioFile(byte_io) as source:
                recognizer.adjust_for_ambient_noise(source, duration=1.5)
                audio = recognizer.record(source)
            query = recognizer.recognize_google(audio, language="ru-RU", show_all=False).lower()
            print(f"  {Y}вы:{RS} {query}")
            log(f"STT: {query}", "OK"); return query
        except sr.UnknownValueError:
            notify("Не расслышал, сэр.", "warn"); return ""
        except sr.RequestError as e:
            notify(f"Google STT: {e}", "err"); return ""
        except Exception as e:
            notify(f"sounddevice: {e} → PyAudio", "warn")
    try:
        print(f"\n  {M}◉ {C}Слушаю...{RS}")
        with sr.Microphone() as source:
            recognizer.adjust_for_ambient_noise(source, duration=1.5)
            recognizer.pause_threshold = 1.5
            audio = recognizer.listen(source, timeout=5, phrase_time_limit=8)
        query = recognizer.recognize_google(audio, language="ru-RU", show_all=False).lower()
        print(f"  {Y}вы:{RS} {query}")
        log(f"STT(mic): {query}", "OK"); return query
    except sr.WaitTimeoutError:
        notify("Тишина.", "info"); return ""
    except sr.UnknownValueError:
        notify("Не расслышал, сэр.", "warn"); return ""
    except sr.RequestError as e:
        notify(f"Google STT: {e}", "err"); return ""
    except Exception as e:
        notify(f"Микрофон: {e}", "err"); return ""

# ══════════════════════════════════════════════════════════════════════
# БЛОК 12 — СТЕГАНОГРАФИЯ (LSB + AES-256)
# ══════════════════════════════════════════════════════════════════════
def _steg_derive_key(password: str) -> bytes:
    if not HAS_STEG:
        raise RuntimeError("Pillow / pycryptodome не установлены")
    return SHA256.new(password.encode("utf-8")).digest()

def steg_encrypt_text(text: str, password: str) -> bytes:
    key = _steg_derive_key(password)
    cipher = AES.new(key, AES.MODE_CBC)
    ct = cipher.encrypt(pad(text.encode("utf-8"), AES.block_size))
    return cipher.iv + ct

def steg_decrypt_text(data: bytes, password: str) -> str:
    key = _steg_derive_key(password)
    iv, ct = data[:16], data[16:]
    cipher = AES.new(key, AES.MODE_CBC, iv=iv)
    return unpad(cipher.decrypt(ct), AES.block_size).decode("utf-8")

def steg_embed(image_path: str, secret_bytes: bytes, out_path: str) -> bool:
    if not HAS_STEG:
        notify("Pillow не установлен", "err"); return False
    try:
        img = Image.open(image_path).convert("RGB")
        pixels = list(img.getdata())

        payload = secret_bytes + STEG_DELIMITER
        payload = struct.pack("<I", len(payload)) + payload

        max_bytes = (len(pixels) * 3) // 8
        if len(payload) > max_bytes:
            notify(f"steg_embed: нужно {len(payload)} байт, доступно {max_bytes}", "warn")
            return False

        bits = []
        for byte in payload:
            for bit in range(7, -1, -1):
                bits.append((byte >> bit) & 1)

        new_pixels = []
        bit_idx = 0
        for px in pixels:
            r, g, b = px
            if bit_idx < len(bits):
                r = (r & 0xFE) | bits[bit_idx]; bit_idx += 1
            if bit_idx < len(bits):
                g = (g & 0xFE) | bits[bit_idx]; bit_idx += 1
            if bit_idx < len(bits):
                b = (b & 0xFE) | bits[bit_idx]; bit_idx += 1
            new_pixels.append((r, g, b))

        result = Image.new("RGB", img.size)
        result.putdata(new_pixels)
        result.save(out_path, format="PNG")
        notify(f"steg_embed: сохранено → {os.path.basename(out_path)}", "ok")
        return True
    except Exception as e:
        notify(f"steg_embed: {e}", "err"); return False

def steg_extract(image_path: str, password: str) -> str | None:
    if not HAS_STEG:
        notify("Pillow не установлен", "err"); return None
    try:
        img = Image.open(image_path).convert("RGB")
        pixels = list(img.getdata())

        bits = []
        for px in pixels:
            for channel in px:
                bits.append(channel & 1)

        def bits_to_bytes(bit_list: list) -> bytes:
            result = bytearray()
            for i in range(0, len(bit_list) - 7, 8):
                byte = 0
                for b in range(8):
                    byte = (byte << 1) | bit_list[i + b]
                result.append(byte)
            return bytes(result)

        raw = bits_to_bytes(bits)
        if len(raw) < 4:
            return None

        length = struct.unpack("<I", raw[:4])[0]
        if length > len(raw) - 4 or length > 10_000_000:
            return None

        payload = raw[4:4 + length]
        if not payload.endswith(STEG_DELIMITER):
            return None

        secret_bytes = payload[:-len(STEG_DELIMITER)]
        return steg_decrypt_text(secret_bytes, password)
    except Exception as e:
        notify(f"steg_extract: {e}", "warn"); return None

# ══════════════════════════════════════════════════════════════════════
# БЛОК 13 — РОУТЕР КОМАНД
# ══════════════════════════════════════════════════════════════════════
def handle_commands(cmd: str) -> bool:
    cmd = cmd.lower().strip()
    if "лог" in cmd or "история команд" in cmd:
        show_log(); return True
    if "сменить пин" in cmd or "изменить пин" in cmd:
        change_pin(); return True
    if any(x in cmd for x in ["забудь всё","очисти память","сбрось память"]):
        clear_history(); return True
    if "покажи историю" in cmd or "что мы обсуждали" in cmd:
        show_history(); return True
    if "статус" in cmd:
        cpu  = psutil.cpu_percent(interval=0.5)
        ram  = psutil.virtual_memory().percent
        disk = psutil.disk_usage("C:\\").percent
        print(f"\n{_hline()}")
        print_progress("CPU    ", cpu)
        print_progress("RAM    ", ram)
        print_progress("Диск C:", disk)
        h = int(time.time()-psutil.boot_time())//3600
        m = (int(time.time()-psutil.boot_time())%3600)//60
        print(_sep())
        print(f"  {DIM}аптайм {h}ч {m}м  ·  процессов {len(psutil.pids())}{RS}")
        print(_hline() + "\n")
        if app: app.update_stats_force()
        speak(f"ЦП: {cpu}%, ОЗУ: {ram}%. Системы в норме, сэр.")
        return True
    if "тест связи" in cmd:
        speak("Анализирую канал, сэр.")
        try:
            res = subprocess.check_output("ping -n 1 google.com", shell=True).decode('cp866')
            p   = res.split("Среднее =")[1].strip() if "Среднее =" in res else "—"
            speak(f"Пинг: {p}. Связь стабильна, сэр.")
        except: speak("Сэр, ошибка пинга.")
        return True
    if any(x in cmd for x in ["время","час"]):
        speak(f"Сейчас {datetime.now().strftime('%H:%M')}, сэр."); return True
    if "очист" in cmd:
        if app: app.clear_log()
        speak("Терминал очищен, сэр."); return True
    if any(x in cmd for x in ["презентация","создай презентацию","создать презентацию","сделай презентацию"]):
        if app: app.after(0, app._pptx_dialog)
        else: create_pptx_terminal()
        return True
    if any(x in cmd for x in ["создай документ","создать документ","новый документ"]):
        if app: app.after(0, app._doc_dialog)
        else: create_doc_terminal()
        return True
    app_map = {
        ("телеграм","telegram","тг"):
            lambda: (os.startfile(r"C:\Users\Aruna\Desktop\Telegram Desktop.lnk")
                     if os.path.exists(r"C:\Users\Aruna\Desktop\Telegram Desktop.lnk")
                     else webbrowser.open("https://web.telegram.org")),
        ("пинтерест","pinterest"):  lambda: webbrowser.open("https://pinterest.com"),
        ("гугл","google"):          lambda: webbrowser.open("https://google.com"),
        ("ютуб","youtube"):         lambda: webbrowser.open("https://youtube.com"),
        ("инстаграм","instagram","инст"): lambda: (
            os.startfile(r"D:\Instagram.lnk") if os.path.exists(r"D:\Instagram.lnk")
            else webbrowser.open("https://instagram.com")),
        ("калькулятор",):  lambda: subprocess.Popen("calc.exe"),
        ("проводник",):    lambda: subprocess.Popen("explorer.exe"),
        ("музыка","яндекс музыку"): lambda: (
            os.startfile(r"C:\Users\Aruna\AppData\Local\Programs\YandexMusic\Яндекс Музыка.exe")
            if os.path.exists(r"C:\Users\Aruna\AppData\Local\Programs\YandexMusic\Яндекс Музыка.exe")
            else webbrowser.open("https://music.yandex.ru")),
        ("стим","steam"): lambda: os.system("start steam://open/main"),
        ("обои",):         lambda: os.system("start steam://rungameid/431960"),
    }
    for keywords, action in app_map.items():
        if any(k in cmd for k in keywords):
            speak(f"Открываю {keywords[0]}, сэр.")
            try: action()
            except Exception as e: speak(f"Сэр, ошибка: {e}")
            return True
    for name, url in STEAM_GAMES.items():
        if name in cmd:
            speak(f"Запускаю {name}. Удачной охоты, сэр.")
            os.system(f"start {url}"); return True
    if any(x in cmd for x in ["майнкрафт","minecraft"]):
        mc = r"D:\LL.exe"
        if os.path.exists(mc): os.startfile(mc); speak("Запускаю Майнкрафт, сэр.")
        else: speak("Сэр, лаунчер не найден.")
        return True
    if any(x in cmd for x in ["курс","валют","доллар","рубл","сум"]):
        rates = get_currency_rates()
        if rates: speak(f"Доллар — {int(rates['usd_uzs'])} сум. Рубль — {rates['rub_uzs']} сума, сэр.")
        else: speak("Сэр, не удалось получить курс.")
        return True
    if any(x in cmd for x in ["погода","температура","на улице"]):
        w = get_weather()
        if w: speak(f"Погода в Фергане, сэр. {w['desc']}. Температура {w['temp']} градусов.")
        else: speak("Сэр, погода недоступна.")
        return True
    triggers = ["нарисуй","сгенерируй","создай картинку","генерируй изображение"]
    if any(x in cmd for x in triggers):
        prompt = cmd
        for t in triggers: prompt = prompt.replace(t,"").strip()
        if prompt: threading.Thread(target=generate_image, args=(prompt,), daemon=True).start()
        else: speak("Сэр, скажите что нарисовать.")
        return True
    return False

# ══════════════════════════════════════════════════════════════════════
# БЛОК 14 — ТЕРМИНАЛЬНЫЕ ДИАЛОГИ
# ══════════════════════════════════════════════════════════════════════
def create_doc_terminal():
    speak("Введите название документа, сэр.")
    title = input(f"  {Y}Название: {W}").strip()
    if not title: return
    print(f"  {Y}Текст (пустая строка = конец):{RS}")
    lines = []
    while True:
        line = input(f"  {DIM}>{RS} ")
        if not line.strip(): break
        lines.append(line)
    if not lines: return
    filepath = _create_docx(title, "\n".join(lines))
    if filepath:
        speak(f"Документ {title} создан, сэр.")
        try: os.startfile(filepath)
        except: pass
    else: speak("Сэр, ошибка создания документа.")

def create_pptx_terminal():
    speak("Введите тему презентации, сэр.")
    topic = input(f"  {Y}Тема: {W}").strip()
    if not topic: return
    while True:
        try:
            n = int(input(f"  {Y}Слайдов (3-10): {W}").strip())
            if 3 <= n <= 10: break
        except: pass
    filepath = generate_pptx(topic, n)
    if filepath:
        speak("Презентация готова, сэр.")
        try: os.startfile(filepath)
        except: pass

# ══════════════════════════════════════════════════════════════════════
# БЛОК 15 — ГЕНЕРАЦИЯ ИЗОБРАЖЕНИЙ
# ══════════════════════════════════════════════════════════════════════
def _translate_to_english(prompt: str) -> str:
    try:
        import g4f
        translated = str(g4f.ChatCompletion.create(
            model=g4f.models.gpt_4,
            messages=[{"role": "user", "content":
                f"Translate to English for image generation. Return ONLY translation, no explanation: {prompt}"}],
            stream=False
        )).strip()
        if translated and len(translated) > 3 and translated != prompt:
            return translated
    except Exception as e:
        notify(f"Перевод: {e}", "warn")
    return prompt

def _try_pollinations(prompt_en: str, filepath: str) -> bool:
    try:
        import urllib.parse
        encoded = urllib.parse.quote(prompt_en)
        seed = random.randint(1, 99999)
        url = (
            f"https://image.pollinations.ai/prompt/{encoded}"
            f"?width=1024&height=1024&seed={seed}&nologo=true&enhance=true"
        )
        notify("Pollinations.ai → генерирую...", "sys")
        resp = requests.get(url, timeout=60, stream=True)
        if resp.status_code == 200 and len(resp.content) > 5000:
            with open(filepath, "wb") as f:
                f.write(resp.content)
            return True
        else:
            notify(f"Pollinations: {resp.status_code} / размер {len(resp.content)}", "warn")
    except Exception as e:
        notify(f"Pollinations: {e}", "warn")
    return False

def _try_hf_inference(prompt_en: str, filepath: str) -> bool:
    headers = {
        "Authorization": f"Bearer {HF_TOKEN}",
        "Content-Type": "application/json",
        "x-wait-for-model": "true",
    }
    models = [
        "black-forest-labs/FLUX.1-schnell",
        "stabilityai/stable-diffusion-3-medium-diffusers",
        "stabilityai/stable-diffusion-xl-base-1.0",
    ]
    payload = {
        "inputs": prompt_en,
        "parameters": {"num_inference_steps": 20, "guidance_scale": 7.0}
    }
    for model in models:
        try:
            url = f"https://router.huggingface.co/hf-inference/models/{model}"
            notify(f"HF → {model.split('/')[-1]}...", "sys")
            resp = requests.post(url, headers=headers, json=payload, timeout=90)
            if resp.status_code == 200 and len(resp.content) > 5000:
                with open(filepath, "wb") as f:
                    f.write(resp.content)
                return True
            else:
                notify(f"{model.split('/')[-1]}: {resp.status_code}", "warn")
        except requests.exceptions.Timeout:
            notify(f"{model.split('/')[-1]}: таймаут", "warn")
        except Exception as e:
            notify(f"{model.split('/')[-1]}: {e}", "warn")
    return False

def _try_g4f_image(prompt_en: str, filepath: str) -> bool:
    try:
        import g4f
        from g4f.client import Client
        import base64
        client = Client()
        response = client.images.generate(
            model="flux", prompt=prompt_en, response_format="b64_json"
        )
        if response.data and response.data[0].b64_json:
            img_data = base64.b64decode(response.data[0].b64_json)
            with open(filepath, "wb") as f:
                f.write(img_data)
            return True
    except Exception as e:
        notify(f"g4f image: {e}", "warn")
    return False

def generate_image(prompt: str):
    speak("Генерирую изображение, сэр. Это займёт несколько секунд.")
    notify(f"Промпт: {prompt}", "sys")
    os.makedirs(JARVIS_DIR, exist_ok=True)
    filename = f"geralt_img_{datetime.now().strftime('%d%m%Y_%H%M%S')}.png"
    filepath = os.path.join(JARVIS_DIR, filename)
    prompt_en = _translate_to_english(prompt)
    if prompt_en != prompt:
        notify(f"EN промпт: {prompt_en}", "sys")
    success = (
        _try_pollinations(prompt_en, filepath) or
        _try_hf_inference(prompt_en, filepath) or
        _try_g4f_image(prompt_en, filepath)
    )
    if success and os.path.exists(filepath) and os.path.getsize(filepath) > 1000:
        notify(f"✓ Изображение сохранено → {filename}", "ok")
        log(f"Изображение: {filename}", "OK")
        try: os.startfile(filepath)
        except Exception as e: notify(f"Открытие: {e}", "warn")
        try:
            with open(filepath, "rb") as photo:
                bot.send_photo(OWNER_ID, photo,
                    caption=(
                        f"🎨 *Готово, сэр!*\n"
                        f"📝 `{prompt}`\n"
                        f"🔤 EN: `{prompt_en[:150]}`"
                    ),
                    parse_mode="Markdown")
            speak("Картинка готова и отправлена в Telegram, сэр.")
        except Exception as e:
            speak("Картинка готова и открыта, сэр.")
            log(f"TG send photo: {e}", "WARN")
        if app:
            app.write_log(f"✓ Изображение: {filename}", tag="ok")
    else:
        speak("Сэр, не удалось сгенерировать изображение. Проверьте интернет-соединение.")
        log("generate_image: все методы провалились", "WARN")
        if app:
            app.write_log("✗ Генерация не удалась. Проверьте интернет.", tag="err")

# ══════════════════════════════════════════════════════════════════════
# БЛОК 16 — GPT / ПАМЯТЬ
# ══════════════════════════════════════════════════════════════════════
def ask_gpt(cmd: str):
    global chat_history
    import g4f
    chat_history.append({"role":"user","content":cmd})
    if len(chat_history) > MAX_HISTORY:
        chat_history = chat_history[-MAX_HISTORY:]
    messages = [
        {"role":"system","content":SYSTEM_PROMPT},
        {"role":"system","content":"НАПОМИНАНИЕ: Обращайся к пользователю ТОЛЬКО 'сэр'."},
    ] + chat_history
    try:
        res    = g4f.ChatCompletion.create(model=g4f.models.gpt_4, messages=messages, stream=False)
        answer = _sanitize_response(str(res).strip())
        chat_history.append({"role":"assistant","content":answer})
        speak(answer); log(f"GPT: {cmd[:50]}", "OK")
    except Exception as e:
        speak("Сэр, ошибка ядра ИИ."); log(f"GPT: {e}", "WARN")

def clear_history():
    global chat_history
    chat_history = []
    speak("Память очищена, сэр."); log("История очищена", "OK")

def show_history():
    if not chat_history: notify("История пуста.", "warn"); return
    print(f"\n{_hline()}")
    for i, msg in enumerate(chat_history):
        tag = Y if msg["role"]=="user" else M
        who = "вы" if msg["role"]=="user" else "геральт"
        print(f"  {tag}[{i+1}] {who:<8}{RS} {msg['content'][:70]}")
    print(_hline() + "\n")

# ══════════════════════════════════════════════════════════════════════
# БЛОК 17 — РЕЖИМ
# ══════════════════════════════════════════════════════════════════════
def try_switch_mode(cmd: str) -> bool:
    global INTERACTION_MODE
    if any(t in cmd for t in ["/голос","переключись на голос","голосовой режим"]):
        INTERACTION_MODE = "voice"
        if app: app.mode_btn.set("voice")
        speak("Переключаюсь на голосовой режим, сэр."); return True
    if any(t in cmd for t in ["/чат","переключись на чат","режим чата"]):
        INTERACTION_MODE = "chat"
        if app: app.mode_btn.set("chat")
        speak("Переключаюсь в режим чата, сэр."); return True
    return False

def get_input() -> str:
    if INTERACTION_MODE == "chat":
        try: return input(f"  {M}geralt{DIM}@system{RS} {M}⚔{RS} ").strip().lower()
        except (EOFError, KeyboardInterrupt): return ""
    else: return listen()

# ══════════════════════════════════════════════════════════════════════
# БЛОК 18 — GUI v5.2  (FIX: _voice_loop переработан)
# ══════════════════════════════════════════════════════════════════════
class GeraltUI(ctk.CTk):
    def __init__(self):
        super().__init__()
        self._theme_name = "dark"
        self.T = THEMES["dark"]
        self.title("GERALT SYSTEM  ⚔  v5.2")
        self.geometry("1280x800")
        self.minsize(960, 620)
        ctk.set_appearance_mode("dark")
        ctk.set_default_color_theme("blue")
        self._sidebar_collapsed = False
        self._cmd_count = 0
        self._session_start = time.time()
        self._voice_active = True
        self._apply_theme_bg()
        self._build_all()
        self._update_clock()
        self._update_stats()
        self._animate_status_dot()
        threading.Thread(target=self._voice_loop, daemon=True).start()

    def _voice_loop(self):
        """
        FIX v5.2 — голосовой цикл.
        Ждём SPEAK_LOCK (он сброшен пока идёт TTS) + доп. задержку,
        только потом вызываем listen().
        """
        while self._voice_active:
            if INTERACTION_MODE != "voice":
                time.sleep(0.3)
                continue

            # Ждём пока TTS закончит (SPEAK_LOCK.set() вызывается в _tts_thread ПОСЛЕ паузы)
            acquired = SPEAK_LOCK.wait(timeout=30)
            if not acquired:
                # Таймаут — проверяем флаг и продолжаем
                continue

            if not self._voice_active:
                break

            # Дополнительная страховочная задержка перед включением микрофона
            time.sleep(TTS_LISTEN_DELAY)

            # Ещё раз проверяем: вдруг за время задержки началась новая озвучка
            if IS_SPEAKING or not self._voice_active:
                continue

            try:
                query = listen()
                if query and query.strip():
                    self.write_log(f"Вы: {query}", tag="user")
                    threading.Thread(
                        target=self._process_cmd, args=(query,), daemon=True
                    ).start()
            except Exception as e:
                notify(f"Voice loop: {e}", "err")
                time.sleep(1)

    def destroy(self):
        self._voice_active = False
        SPEAK_LOCK.set()   # разблокируем поток чтобы он мог выйти
        super().destroy()

    def _apply_theme_bg(self):
        self.configure(fg_color=self.T["bg"])

    def _switch_theme(self):
        self._theme_name = "light" if self._theme_name == "dark" else "dark"
        self.T = THEMES[self._theme_name]
        ctk.set_appearance_mode("light" if self._theme_name == "light" else "dark")
        for widget in self.winfo_children(): widget.destroy()
        self._apply_theme_bg()
        self._build_all()
        self._update_clock()
        self._update_stats()
        self._animate_status_dot()
        notify(f"Тема переключена: {self._theme_name}", "sys")

    def _build_all(self):
        self._build_titlebar()
        body = ctk.CTkFrame(self, fg_color=self.T["bg"], corner_radius=0)
        body.pack(fill="both", expand=True)
        body.grid_rowconfigure(0, weight=1)
        body.grid_columnconfigure(1, weight=1)
        self._sidebar_frame = self._build_sidebar(body)
        self._sidebar_frame.grid(row=0, column=0, sticky="nsew")
        main = ctk.CTkFrame(body, fg_color=self.T["bg"], corner_radius=0)
        main.grid(row=0, column=1, sticky="nsew")
        main.grid_rowconfigure(0, weight=1)
        main.grid_columnconfigure(0, weight=1)
        self._build_main(main)

    def _build_titlebar(self):
        bar = ctk.CTkFrame(self, height=52, fg_color=self.T["surface"], corner_radius=0)
        bar.pack(fill="x"); bar.pack_propagate(False)
        ctk.CTkFrame(bar, height=2, fg_color=self.T["accent"], corner_radius=0).pack(fill="x", side="top")
        inner = ctk.CTkFrame(bar, fg_color="transparent")
        inner.pack(fill="both", expand=True, padx=16)
        left = ctk.CTkFrame(inner, fg_color="transparent")
        left.pack(side="left", fill="y")
        ctk.CTkButton(left, text="☰", width=32, height=32, fg_color="transparent",
                      hover_color=self.T["surface2"], text_color=self.T["text2"],
                      font=("Consolas",16), command=self._toggle_sidebar
                      ).pack(side="left", padx=(0,12), pady=10)
        logo = ctk.CTkFrame(left, fg_color=self.T["accent_soft"], corner_radius=8, width=32, height=32)
        logo.pack(side="left", pady=10); logo.pack_propagate(False)
        ctk.CTkLabel(logo, text="⚔", font=("Consolas",16),
                     text_color=self.T["accent"]).place(relx=0.5, rely=0.5, anchor="center")
        ctk.CTkLabel(left, text="GERALT", font=("Consolas",17,"bold"),
                     text_color=self.T["text"]).pack(side="left", padx=(10,0))
        ctk.CTkLabel(left, text="AI", font=("Consolas",13),
                     text_color=self.T["accent"]).pack(side="left", padx=(4,0))
        ctk.CTkLabel(left, text=" v5.2", font=("Consolas",11),
                     text_color=self.T["text3"]).pack(side="left")
        mid = ctk.CTkFrame(inner, fg_color="transparent")
        mid.pack(side="left", padx=28, fill="y")
        self._status_dot = ctk.CTkLabel(mid, text="●", font=("Consolas",10),
                                        text_color=self.T["success"])
        self._status_dot.pack(side="left")
        ctk.CTkLabel(mid, text=" ONLINE", font=("Consolas",10),
                     text_color=self.T["text3"]).pack(side="left")
        ctk.CTkFrame(mid, width=1, height=16, fg_color=self.T["border"]).pack(side="left", padx=10)
        audio_col = self.T["success"] if HAS_VOICE else self.T["danger"]
        ctk.CTkLabel(mid, text=f"AUDIO:{AUDIO_BACKEND.upper()}", font=("Consolas",9),
                     text_color=audio_col).pack(side="left")
        ctk.CTkFrame(mid, width=1, height=16, fg_color=self.T["border"]).pack(side="left", padx=10)
        mic_col = self.T["success"] if (HAS_MIC or HAS_SR) else self.T["danger"]
        ctk.CTkLabel(mid, text=f"MIC:{MIC_BACKEND.upper()}", font=("Consolas",9),
                     text_color=mic_col).pack(side="left")
        ctk.CTkFrame(mid, width=1, height=16, fg_color=self.T["border"]).pack(side="left", padx=10)
        steg_col = self.T["success"] if HAS_STEG else self.T["warn"]
        ctk.CTkLabel(mid, text=f"STEG:{'ON' if HAS_STEG else 'OFF'}", font=("Consolas",9),
                     text_color=steg_col).pack(side="left")
        right = ctk.CTkFrame(inner, fg_color="transparent")
        right.pack(side="right", fill="y")
        self.clock_lbl = ctk.CTkLabel(right, text="", font=("Consolas",12),
                                      text_color=self.T["text3"])
        self.clock_lbl.pack(side="right", padx=(14,0))
        theme_icon = "☀" if self._theme_name == "dark" else "☾"
        theme_tip  = "Светлая тема" if self._theme_name == "dark" else "Тёмная тема"
        self._theme_btn = ctk.CTkButton(
            right, text=f"{theme_icon}  {theme_tip}", width=130, height=30,
            fg_color=self.T["surface2"], hover_color=self.T["border2"],
            border_width=1, border_color=self.T["border"],
            text_color=self.T["text2"], font=("Consolas",11), corner_radius=8,
            command=self._switch_theme)
        self._theme_btn.pack(side="right", pady=10, padx=(8,0))
        self.mode_btn = ctk.CTkSegmentedButton(
            right, values=["chat","voice"], command=self._on_mode_switch,
            font=("Consolas",11,"bold"),
            selected_color=self.T["accent"], selected_hover_color=self.T["accent_hover"],
            unselected_color=self.T["surface2"], unselected_hover_color=self.T["border"],
            fg_color=self.T["surface2"], text_color=self.T["text"], width=150, height=30)
        self.mode_btn.pack(side="right", pady=10)
        self.mode_btn.set("voice")
        ctk.CTkFrame(bar, height=1, fg_color=self.T["border"], corner_radius=0).pack(fill="x", side="bottom")

    def _build_sidebar(self, parent):
        side = ctk.CTkFrame(parent, width=230, fg_color=self.T["surface"], corner_radius=0)
        side.pack_propagate(False)
        ctk.CTkFrame(side, width=1, fg_color=self.T["border"], corner_radius=0).pack(side="right", fill="y")
        content = ctk.CTkScrollableFrame(side, fg_color="transparent",
                                          scrollbar_button_color=self.T["border"],
                                          scrollbar_button_hover_color=self.T["border2"])
        content.pack(fill="both", expand=True, pady=8)
        self._sb_section(content, "БЫСТРЫЕ КОМАНДЫ")
        quick = [
            ("🎮","Ведьмак",        lambda: self._quick("запусти ведьмак")),
            ("🎨","Генерация фото", self._image_dialog),
            ("🌐","Google",         lambda: self._quick("гугл")),
            ("📺","YouTube",        lambda: self._quick("ютуб")),
            ("💱","Курс валют",     lambda: self._quick("курс валют")),
            ("🌤","Погода",         lambda: self._quick("погода")),
            ("💻","Статус ПК",      lambda: self._quick("статус")),
            ("📊","Презентация",    lambda: self._quick("создай презентацию")),
            ("📄","Документ",       lambda: self._quick("создай документ")),
        ]
        for icon, label, fn in quick:
            self._sb_btn(content, icon, label, fn)
        self._sb_section(content, "СТЕГАНОГРАФИЯ")
        self._sb_btn(content, "🔒", "Скрыть текст в фото", self._steg_embed_dialog)
        self._sb_btn(content, "🔓", "Извлечь текст",       self._steg_extract_dialog)
        self._sb_section(content, "СИСТЕМА")
        self._sb_btn(content, "🗑", "Очистить лог",    self.clear_log)
        self._sb_btn(content, "🔄", "Сбросить память", lambda: self._quick("забудь всё"))
        self._sb_btn(content, "📋", "История",          lambda: self._quick("покажи историю"))
        self._sb_btn(content, "🔴", "Выключить",        lambda: self._quick("отключись"), danger=True)
        self._sb_section(content, "МОНИТОРИНГ")
        mon = ctk.CTkFrame(content, fg_color="transparent")
        mon.pack(fill="x", padx=12, pady=4)
        self.cpu_bar  = self._monitor_row(mon, "CPU",  self.T["accent2"])
        self.ram_bar  = self._monitor_row(mon, "RAM",  self.T["accent"])
        self.disk_bar = self._monitor_row(mon, "DISK", self.T["warn"])
        return side

    def _sb_section(self, parent, text: str):
        f = ctk.CTkFrame(parent, fg_color="transparent")
        f.pack(fill="x", padx=14, pady=(14,4))
        ctk.CTkLabel(f, text=text, font=("Consolas",9,"bold"),
                     text_color=self.T["text3"]).pack(side="left")
        ctk.CTkFrame(f, height=1, fg_color=self.T["border"]).pack(
            side="left", fill="x", expand=True, padx=(8,0), pady=6)

    def _sb_btn(self, parent, icon: str, label: str, cmd, danger=False):
        color    = self.T["danger"]      if danger else self.T["text2"]
        hover_bg = self.T["danger_soft"] if danger else self.T["surface2"]
        act_bg   = self.T["danger_soft"] if danger else self.T["accent_soft"]
        btn_f    = ctk.CTkFrame(parent, fg_color="transparent", cursor="hand2")
        btn_f.pack(fill="x", padx=8, pady=1)
        inner = ctk.CTkFrame(btn_f, fg_color="transparent", corner_radius=8)
        inner.pack(fill="x", padx=2)
        icon_lbl = ctk.CTkLabel(inner, text=icon, font=("Consolas",13),
                                 text_color=color, width=28)
        icon_lbl.pack(side="left", padx=(10,0), pady=7)
        text_lbl = ctk.CTkLabel(inner, text=label, font=("Consolas",12),
                                 text_color=color, anchor="w")
        text_lbl.pack(side="left", padx=8, pady=7, fill="x", expand=True)
        def on_enter(e): inner.configure(fg_color=hover_bg)
        def on_leave(e): inner.configure(fg_color="transparent")
        def on_click(e):
            inner.configure(fg_color=act_bg)
            self.after(120, lambda: inner.configure(fg_color=hover_bg))
            cmd()
        for w in [inner, icon_lbl, text_lbl]:
            w.bind("<Enter>", on_enter); w.bind("<Leave>", on_leave); w.bind("<Button-1>", on_click)

    def _monitor_row(self, parent, label: str, color: str):
        row = ctk.CTkFrame(parent, fg_color=self.T["surface2"], corner_radius=8)
        row.pack(fill="x", pady=3)
        top = ctk.CTkFrame(row, fg_color="transparent")
        top.pack(fill="x", padx=10, pady=(8,2))
        ctk.CTkLabel(top, text=label, font=("Consolas",10,"bold"),
                     text_color=self.T["text3"], anchor="w").pack(side="left")
        val_lbl = ctk.CTkLabel(top, text="0%", font=("Consolas",10,"bold"),
                                text_color=color, anchor="e")
        val_lbl.pack(side="right")
        bar = ctk.CTkProgressBar(row, height=4, corner_radius=2,
                                  fg_color=self.T["border"], progress_color=color)
        bar.pack(fill="x", padx=10, pady=(0,8)); bar.set(0)
        return bar, val_lbl

    def _toggle_sidebar(self):
        self._sidebar_collapsed = not self._sidebar_collapsed
        self._sidebar_frame.configure(width=0 if self._sidebar_collapsed else 230)

    def _build_main(self, parent):
        stats_row = ctk.CTkFrame(parent, fg_color="transparent")
        stats_row.pack(fill="x", padx=16, pady=(12,0))
        self._stat_labels = {}
        stat_items = [
            ("⚔","Сессия",  self._get_session_time, self.T["accent"]),
            ("◈","Команды", self._get_cmd_count,     self.T["accent2"]),
            ("⊙","Голос",   lambda: ("ON" if HAS_VOICE else "OFF"), self.T["success"]),
            ("🔒","Стег",   lambda: ("ON" if HAS_STEG else "OFF"),  self.T["warn"]),
        ]
        for icon, name, getter, color in stat_items:
            card = ctk.CTkFrame(stats_row, fg_color=self.T["surface"],
                                corner_radius=10, border_width=1,
                                border_color=self.T["border"])
            card.pack(side="left", padx=(0,8), fill="y")
            ctk.CTkLabel(card, text=f"{icon}  {name}", font=("Consolas",9),
                         text_color=self.T["text3"]).pack(padx=14, pady=(8,0), anchor="w")
            lbl = ctk.CTkLabel(card, text=getter(), font=("Consolas",14,"bold"),
                               text_color=color)
            lbl.pack(padx=14, pady=(0,8), anchor="w")
            self._stat_labels[name] = (lbl, getter)
        log_outer = ctk.CTkFrame(parent, fg_color=self.T["surface"],
                                  corner_radius=12, border_width=1,
                                  border_color=self.T["border"])
        log_outer.pack(fill="both", expand=True, padx=16, pady=12)
        log_hdr = ctk.CTkFrame(log_outer, fg_color="transparent", height=42)
        log_hdr.pack(fill="x", padx=16, pady=(10,0)); log_hdr.pack_propagate(False)
        ctk.CTkLabel(log_hdr, text="▸", font=("Consolas",14),
                     text_color=self.T["accent"]).pack(side="left")
        ctk.CTkLabel(log_hdr, text=" TERMINAL", font=("Consolas",12,"bold"),
                     text_color=self.T["text"]).pack(side="left")
        btn_f = ctk.CTkFrame(log_hdr, fg_color="transparent")
        btn_f.pack(side="right")
        self.msg_count_lbl = ctk.CTkLabel(btn_f, text="0 строк",
                                           font=("Consolas",10), text_color=self.T["text3"])
        self.msg_count_lbl.pack(side="right", padx=(8,0))
        ctk.CTkButton(btn_f, text="очистить", width=70, height=24,
                      fg_color=self.T["surface2"], hover_color=self.T["border"],
                      border_width=1, border_color=self.T["border"],
                      text_color=self.T["text3"], font=("Consolas",10),
                      corner_radius=6, command=self.clear_log).pack(side="right")
        ctk.CTkFrame(log_outer, height=1, fg_color=self.T["border"]).pack(fill="x", pady=(8,0))
        self.log_box = ctk.CTkTextbox(
            log_outer, font=("Consolas",13), fg_color="transparent",
            text_color=self.T["text"], wrap="word",
            scrollbar_button_color=self.T["border"],
            scrollbar_button_hover_color=self.T["border2"])
        self.log_box.pack(fill="both", expand=True, padx=8, pady=(0,6))
        tb = self.log_box._textbox
        tb.tag_configure("user",   foreground=self.T["accent2"])
        tb.tag_configure("geralt", foreground=self.T["accent"])
        tb.tag_configure("sys",    foreground=self.T["text3"])
        tb.tag_configure("err",    foreground=self.T["danger"])
        tb.tag_configure("ok",     foreground=self.T["success"])
        tb.tag_configure("warn",   foreground=self.T["warn"])
        tb.tag_configure("time",   foreground=self.T["text3"])
        tb.tag_configure("arrow",  foreground=self.T["border2"])
        if not HAS_VOICE:
            self.write_log("⚠ АУДИО НЕДОСТУПНО — pip install playsound==1.2.2", tag="err")
        if not HAS_MIC and not HAS_SR:
            self.write_log("⚠ МИКРОФОН НЕДОСТУПЕН — pip install pyaudio SpeechRecognition", tag="err")
        if not HAS_STEG:
            self.write_log("⚠ СТЕГАНОГРАФИЯ НЕДОСТУПНА — pip install Pillow pycryptodome", tag="warn")
        self._listen_indicator = ctk.CTkLabel(
            log_outer, text="", font=("Consolas", 10),
            text_color=self.T["success"])
        self._listen_indicator.pack(anchor="w", padx=12, pady=(0, 4))
        self._blink_listen_indicator()
        input_outer = ctk.CTkFrame(parent, fg_color=self.T["surface"],
                                    corner_radius=12, border_width=1,
                                    border_color=self.T["border"])
        input_outer.pack(fill="x", padx=16, pady=(0,14))
        input_inner = ctk.CTkFrame(input_outer, fg_color="transparent")
        input_inner.pack(fill="x", padx=12, pady=10)
        self._mode_pill = ctk.CTkLabel(
            input_inner, text="  VOICE  ", font=("Consolas",9,"bold"),
            fg_color=self.T["accent_soft"], text_color=self.T["accent"],
            corner_radius=4, width=56, height=22)
        self._mode_pill.pack(side="left", padx=(0,10))
        self.entry = ctk.CTkEntry(
            input_inner, height=40,
            placeholder_text="Введите команду, сэр...",
            font=("Consolas",13), fg_color=self.T["surface2"],
            border_color=self.T["border"], border_width=1,
            text_color=self.T["text"],
            placeholder_text_color=self.T["text3"], corner_radius=8)
        self.entry.pack(side="left", fill="x", expand=True, padx=(0,10))
        self.entry.bind("<Return>", self._on_enter)
        self.entry.bind("<FocusIn>",  lambda e: self.entry.configure(border_color=self.T["accent"]))
        self.entry.bind("<FocusOut>", lambda e: self.entry.configure(border_color=self.T["border"]))
        ctk.CTkButton(
            input_inner, text="⚔  Отправить", width=140, height=40,
            fg_color=self.T["accent"], hover_color=self.T["accent_hover"],
            font=("Consolas",12,"bold"), text_color=self.T["text_inv"],
            corner_radius=8, command=self._on_enter).pack(side="right")

    def _blink_listen_indicator(self):
        if not self._voice_active: return
        try:
            if INTERACTION_MODE == "voice":
                if IS_SPEAKING:
                    self._listen_indicator.configure(
                        text="  🔇  Геральт говорит... (микрофон заглушен)",
                        text_color=self.T["warn"])
                else:
                    current = self._listen_indicator.cget("text")
                    if "◉" in current:
                        self._listen_indicator.configure(
                            text="  ◯  Слушаю микрофон...", text_color=self.T["text3"])
                    else:
                        self._listen_indicator.configure(
                            text="  ◉  Слушаю микрофон...", text_color=self.T["success"])
            else:
                self._listen_indicator.configure(
                    text="  —  Режим чата (текстовый ввод)", text_color=self.T["text3"])
        except: pass
        self.after(800, self._blink_listen_indicator)

    def _get_session_time(self) -> str:
        elapsed = int(time.time() - self._session_start)
        h, m = elapsed // 3600, (elapsed % 3600) // 60
        return f"{h:02d}:{m:02d}"

    def _get_cmd_count(self) -> str:
        return str(self._cmd_count)

    def _on_mode_switch(self, val):
        global INTERACTION_MODE
        INTERACTION_MODE = val
        pill_color = self.T["accent_soft"]  if val == "voice" else self.T["accent2_soft"]
        text_color = self.T["accent"]       if val == "voice" else self.T["accent2"]
        self._mode_pill.configure(text=f"  {val.upper()}  ",
                                  fg_color=pill_color, text_color=text_color)
        self.write_log(f"Режим: {val}", tag="sys")

    def _quick(self, cmd: str):
        self.write_log(f"Вы: {cmd}", tag="user")
        threading.Thread(target=self._process_cmd, args=(cmd,), daemon=True).start()

    def _on_enter(self, event=None):
        text = self.entry.get().strip()
        if not text: return
        self.entry.delete(0, "end")
        self.write_log(f"Вы: {text}", tag="user")
        threading.Thread(target=self._process_cmd, args=(text,), daemon=True).start()

    def _process_cmd(self, cmd: str):
        self._cmd_count += 1
        cmd_l = cmd.lower()
        if any(x in cmd_l for x in ["отключись","выключись","закройся"]):
            speak("Отключаю систему. До встречи, сэр.")
            log("Система выключена", "OK")
            self.after(1500, self.destroy); return
        if try_switch_mode(cmd_l): return
        if not handle_commands(cmd_l): ask_gpt(cmd)

    def write_log(self, text: str, tag: str = ""):
        def _write():
            now = datetime.now().strftime("%H:%M:%S")
            tb  = self.log_box._textbox
            tb.insert("end", f"[{now}] ", "time")
            tb.insert("end", "▸ ", "arrow")
            _tag = tag
            if not _tag:
                lt = text.lower()
                if lt.startswith("вы:"):         _tag = "user"
                elif lt.startswith("геральт:"):  _tag = "geralt"
                elif any(x in lt for x in ["ошибка","error","✗"]): _tag = "err"
                elif any(x in lt for x in ["✓","готов","успешно"]):  _tag = "ok"
                elif any(x in lt for x in ["⚠","warn"]):             _tag = "warn"
                else:                             _tag = "sys"
            tb.insert("end", f"{text}\n", _tag)
            tb.see("end")
            lines = int(tb.index("end-1c").split(".")[0])
            self.msg_count_lbl.configure(text=f"{lines} строк")
        self.after(0, _write)

    def clear_log(self):
        self.log_box.delete("1.0","end")
        self.msg_count_lbl.configure(text="0 строк")

    def update_stats_force(self): self._update_stats()

    def _update_stats(self):
        try:
            cpu  = psutil.cpu_percent()
            ram  = psutil.virtual_memory().percent
            disk = psutil.disk_usage("C:\\").percent
            def _set(bt, v):
                bar, lbl = bt; bar.set(v/100); lbl.configure(text=f"{v:.0f}%")
            _set(self.cpu_bar, cpu); _set(self.ram_bar, ram); _set(self.disk_bar, disk)
            for name, (lbl, getter) in self._stat_labels.items():
                lbl.configure(text=getter())
        except: pass
        self.after(3000, self._update_stats)

    def _update_clock(self):
        self.clock_lbl.configure(text=datetime.now().strftime("%d.%m.%Y  %H:%M:%S"))
        self.after(1000, self._update_clock)

    def _animate_status_dot(self):
        current    = self._status_dot.cget("text_color")
        next_color = self.T["text3"] if current == self.T["success"] else self.T["success"]
        self._status_dot.configure(text_color=next_color)
        self.after(1200, self._animate_status_dot)

    def _make_dialog(self, title: str, w: int, h: int) -> ctk.CTkToplevel:
        d = ctk.CTkToplevel(self)
        d.title(title); d.geometry(f"{w}x{h}")
        d.configure(fg_color=self.T["surface"])
        d.grab_set(); d.focus(); d.resizable(False, False)
        return d

    def _dlg_header(self, parent, icon: str, title: str):
        f = ctk.CTkFrame(parent, fg_color=self.T["accent_soft"], height=52, corner_radius=0)
        f.pack(fill="x"); f.pack_propagate(False)
        ctk.CTkLabel(f, text=f"{icon}  {title}", font=("Consolas",14,"bold"),
                     text_color=self.T["accent"]).place(relx=0.5, rely=0.5, anchor="center")

    def _dlg_label(self, parent, text: str):
        ctk.CTkLabel(parent, text=text, font=("Consolas",11),
                     text_color=self.T["text2"], anchor="w").pack(anchor="w", padx=24, pady=(10,2))

    def _dlg_entry(self, parent, placeholder: str) -> ctk.CTkEntry:
        e = ctk.CTkEntry(parent, width=430, height=38, placeholder_text=placeholder,
                          font=("Consolas",13), fg_color=self.T["surface2"],
                          border_color=self.T["border"], border_width=1,
                          text_color=self.T["text"], placeholder_text_color=self.T["text3"],
                          corner_radius=8)
        e.pack(padx=24, pady=(0,4))
        e.bind("<FocusIn>",  lambda ev: e.configure(border_color=self.T["accent"]))
        e.bind("<FocusOut>", lambda ev: e.configure(border_color=self.T["border"]))
        return e

    def _dlg_btn(self, parent, label: str, cmd):
        ctk.CTkButton(parent, text=label, height=40, width=430,
                      fg_color=self.T["accent"], hover_color=self.T["accent_hover"],
                      font=("Consolas",13,"bold"), text_color=self.T["text_inv"],
                      corner_radius=8, command=cmd).pack(padx=24, pady=(8,20))

    def _image_dialog(self):
        d = self._make_dialog("Генерация изображения", 480, 220)
        self._dlg_header(d, "🎨", "Генерация изображения")
        self._dlg_label(d, "Опишите что нарисовать:")
        entry = self._dlg_entry(d, "Закат над горами, акварель...")
        entry.focus()
        def _go(event=None):
            p = entry.get().strip(); d.destroy()
            if p:
                self.write_log(f"Вы: нарисуй {p}", tag="user")
                threading.Thread(target=generate_image, args=(p,), daemon=True).start()
        entry.bind("<Return>", _go)
        self._dlg_btn(d, "⚔  Генерировать", _go)

    def _pptx_dialog(self):
        d = self._make_dialog("Создание презентации", 480, 280)
        self._dlg_header(d, "📊", "Создание презентации")
        self._dlg_label(d, "Тема презентации:")
        topic_e = self._dlg_entry(d, "Искусственный интеллект...")
        topic_e.focus()
        self._dlg_label(d, "Количество слайдов (3–10):")
        slides_e = self._dlg_entry(d, "5")
        def _go(event=None):
            topic = topic_e.get().strip()
            try: n = max(3, min(10, int(slides_e.get().strip() or "5")))
            except: n = 5
            d.destroy()
            if topic:
                self.write_log(f"Вы: презентация — {topic} ({n} слайдов)", tag="user")
                speak("Генерирую, сэр.")
                def _run():
                    fp = generate_pptx(topic, n)
                    if fp:
                        speak("Готово, сэр.")
                        try: os.startfile(fp)
                        except: pass
                threading.Thread(target=_run, daemon=True).start()
        slides_e.bind("<Return>", _go)
        self._dlg_btn(d, "⚔  Создать", _go)

    def _doc_dialog(self):
        d = self._make_dialog("Создание документа", 480, 380)
        self._dlg_header(d, "📄", "Документ Word")
        self._dlg_label(d, "Название документа:")
        title_e = self._dlg_entry(d, "Название...")
        title_e.focus()
        self._dlg_label(d, "Текст документа:")
        text_box = ctk.CTkTextbox(d, width=430, height=130, font=("Consolas",12),
                                   fg_color=self.T["surface2"], text_color=self.T["text"],
                                   border_width=1, border_color=self.T["border"], corner_radius=8)
        text_box.pack(padx=24, pady=(0,4))
        def _go():
            title   = title_e.get().strip()
            content = text_box.get("1.0","end").strip()
            d.destroy()
            if not title or not content:
                speak("Сэр, укажите название и текст."); return
            self.write_log(f"Вы: документ — {title}", tag="user")
            def _run():
                fp = _create_docx(title, content)
                if fp:
                    speak(f"Документ {title} создан, сэр.")
                    try: os.startfile(fp)
                    except: pass
                else: speak("Сэр, ошибка.")
            threading.Thread(target=_run, daemon=True).start()
        self._dlg_btn(d, "⚔  Создать", _go)

    # ── СТЕГАНОГРАФИЯ — GUI ДИАЛОГИ ──────────────────────────────────
    def _steg_embed_dialog(self):
        if not HAS_STEG:
            self.write_log("⚠ Установи: pip install Pillow pycryptodome", tag="err")
            speak("Сэр, библиотеки стеганографии не установлены.")
            return
        from tkinter import filedialog
        image_path = filedialog.askopenfilename(
            title="Выберите изображение",
            filetypes=[("Изображения", "*.png *.jpg *.jpeg *.bmp *.tiff"), ("Все файлы", "*.*")]
        )
        if not image_path: return
        d = self._make_dialog("Скрыть текст в фото", 480, 420)
        self._dlg_header(d, "🔒", "Скрыть текст в фото")
        ctk.CTkLabel(d, text=f"Файл: {os.path.basename(image_path)}",
                     font=("Consolas",10), text_color=self.T["text3"]).pack(padx=24, pady=(8,0), anchor="w")
        self._dlg_label(d, "Текст для скрытия:")
        text_box = ctk.CTkTextbox(d, width=430, height=130, font=("Consolas",12),
                                   fg_color=self.T["surface2"], text_color=self.T["text"],
                                   border_width=1, border_color=self.T["border"], corner_radius=8)
        text_box.pack(padx=24, pady=(0,4))
        self._dlg_label(d, "Пароль:")
        pass_e = self._dlg_entry(d, "Секретный пароль...")
        def _go():
            secret = text_box.get("1.0","end").strip()
            password = pass_e.get().strip()
            d.destroy()
            if not secret or not password:
                speak("Сэр, заполните все поля."); return
            self.write_log("🔒 Шифрую и встраиваю...", tag="sys")
            def _run():
                try:
                    encrypted = steg_encrypt_text(secret, password)
                    os.makedirs(JARVIS_DIR, exist_ok=True)
                    out_name = f"steg_{datetime.now().strftime('%d%m%Y_%H%M%S')}.png"
                    out_path = os.path.join(JARVIS_DIR, out_name)
                    ok = steg_embed(image_path, encrypted, out_path)
                    if ok:
                        speak("Готово, сэр. Файл сохранён.")
                        self.write_log(f"✓ Сохранено: {out_path}", tag="ok")
                        try: os.startfile(JARVIS_DIR)
                        except: pass
                    else:
                        speak("Сэр, фото слишком маленькое или ошибка.")
                        self.write_log("✗ Не удалось встроить — фото слишком маленькое?", tag="err")
                except Exception as e:
                    speak("Сэр, ошибка шифрования.")
                    self.write_log(f"✗ steg embed: {e}", tag="err")
            threading.Thread(target=_run, daemon=True).start()
        self._dlg_btn(d, "🔒  Скрыть", _go)

    def _steg_extract_dialog(self):
        if not HAS_STEG:
            self.write_log("⚠ Установи: pip install Pillow pycryptodome", tag="err")
            speak("Сэр, библиотеки стеганографии не установлены.")
            return
        from tkinter import filedialog
        image_path = filedialog.askopenfilename(
            title="Выберите изображение со скрытым текстом",
            filetypes=[("PNG", "*.png"), ("Изображения", "*.png *.jpg *.jpeg *.bmp"), ("Все файлы", "*.*")]
        )
        if not image_path: return
        d = self._make_dialog("Извлечь текст из фото", 480, 250)
        self._dlg_header(d, "🔓", "Извлечь текст из фото")
        ctk.CTkLabel(d, text=f"Файл: {os.path.basename(image_path)}",
                     font=("Consolas",10), text_color=self.T["text3"]).pack(padx=24, pady=(8,0), anchor="w")
        self._dlg_label(d, "Пароль:")
        pass_e = self._dlg_entry(d, "Введите пароль...")
        pass_e.focus()
        def _go(event=None):
            password = pass_e.get().strip()
            d.destroy()
            if not password:
                speak("Сэр, введите пароль."); return
            self.write_log("🔓 Извлекаю текст...", tag="sys")
            def _run():
                try:
                    result = steg_extract(image_path, password)
                    if result is not None:
                        speak("Текст успешно извлечён, сэр.")
                        self.write_log("✓ Скрытое сообщение:", tag="ok")
                        self.after(0, lambda: self._show_steg_result(result))
                    else:
                        speak("Сэр, не удалось извлечь. Неверный пароль или повреждённое фото.")
                        self.write_log("✗ Текст не найден (неверный пароль или не тот файл)", tag="err")
                except Exception as e:
                    speak("Сэр, ошибка расшифровки.")
                    self.write_log(f"✗ steg extract: {e}", tag="err")
            threading.Thread(target=_run, daemon=True).start()
        pass_e.bind("<Return>", _go)
        self._dlg_btn(d, "🔓  Извлечь", _go)

    def _show_steg_result(self, text: str):
        d = self._make_dialog("Скрытое сообщение", 560, 400)
        self._dlg_header(d, "🔓", "Скрытое сообщение")
        box = ctk.CTkTextbox(d, font=("Consolas",13), fg_color=self.T["surface2"],
                              text_color=self.T["success"],
                              border_width=1, border_color=self.T["border"], corner_radius=8)
        box.pack(fill="both", expand=True, padx=24, pady=(12,8))
        box.insert("1.0", text)
        box.configure(state="disabled")
        ctk.CTkButton(d, text="✓  Закрыть", height=36,
                      fg_color=self.T["surface2"], hover_color=self.T["border"],
                      text_color=self.T["text2"], font=("Consolas",12),
                      corner_radius=8, command=d.destroy).pack(padx=24, pady=(0,16))

# ══════════════════════════════════════════════════════════════════════
# БЛОК 19 — ЗАПУСК
# ══════════════════════════════════════════════════════════════════════
def show_logo():
    os.system('cls' if os.name == 'nt' else 'clear')
    now = datetime.now()
    try:    bot.get_me(); tg_status = f"{G}online{RS}"
    except: tg_status = f"{R}offline{RS}"
    voice_status = f"{G}{AUDIO_BACKEND}{RS}" if HAS_VOICE else f"{R}нет{RS}"
    mic_status   = f"{G}{MIC_BACKEND}{RS}"   if HAS_MIC   else (f"{Y}sr-only{RS}" if HAS_SR else f"{R}нет{RS}")
    steg_status  = f"{G}Pillow+AES{RS}"      if HAS_STEG  else f"{Y}не установлен{RS}"
    print()
    print(f"  {M}╔{'═'*54}╗{RS}")
    print(f"  {M}║{RS}  {M}G E R A L T{RS}  {DIM}·{RS}  {W}Personal AI System{RS}  {DIM}·{RS}  {G}v5.2{RS}  {M}  ║{RS}")
    print(f"  {M}╟{'─'*54}╢{RS}")
    print(f"  {M}║{RS}  {Y}tg{RS} {tg_status}   {Y}voice{RS} {voice_status}   {Y}mic{RS} {mic_status}   {M}║{RS}")
    print(f"  {M}║{RS}  {Y}steg{RS} {steg_status}   {DIM}{now.strftime('%d.%m.%Y  %H:%M')}{RS}   {M}║{RS}")
    print(f"  {M}╚{'═'*54}╝{RS}")
    print()
    if not HAS_VOICE:
        print(f"  {R}⚠  ГОЛОС НЕДОСТУПЕН.{RS} Установите:")
        print(f"  {Y}     pip install playsound==1.2.2{RS}\n")
    if not HAS_MIC and not HAS_SR:
        print(f"  {R}⚠  МИКРОФОН НЕДОСТУПЕН.{RS} Установите:")
        print(f"  {Y}     pip install pyaudio SpeechRecognition{RS}\n")
    if not HAS_STEG:
        print(f"  {Y}⚠  СТЕГАНОГРАФИЯ НЕДОСТУПНА.{RS} Установите:")
        print(f"  {Y}     pip install Pillow pycryptodome{RS}\n")

def choose_mode():
    global INTERACTION_MODE
    print(f"{_hline()}")
    print(f"  {C}ВЫБОР РЕЖИМА{RS}")
    print(_sep())
    print(f"  {C}[1]{RS}  Чат-режим")
    print(f"  {C}[2]{RS}  {G}Голосовой режим{RS}  (рекомендуется)")
    print(_sep())
    while True:
        choice = input(f"  {M}geralt{DIM}@system{RS} {M}⚔{RS} ").strip()
        if choice == "1":   INTERACTION_MODE = "chat"; break
        elif choice == "2": INTERACTION_MODE = "voice"; break
        else: notify("Введите 1 или 2.", "warn")
    print(_hline() + "\n")

def main():
    global app
    if not security_check(): sys.exit()
    try:
        from face_auth import authenticate_face
        if not authenticate_face(intruder_callback=notify_intruder):
            log("Face ID не пройден", "BLOCK"); sys.exit()
    except ImportError:
        notify("face_auth не найден — Face ID пропущен.", "warn")
    log("Система запущена", "OK")
    threading.Thread(target=bot.infinity_polling, daemon=True).start()
    notify("Telegram-бот запущен.", "ok")
    threading.Thread(target=scheduler_loop, daemon=True).start()
    show_logo()
    choose_mode()
    app = GeraltUI()
    app.mode_btn.set(INTERACTION_MODE)
    greeting = "Система Геральт активирована. Добро пожаловать домой, сэр."
    app.write_log(greeting, tag="geralt")
    threading.Thread(target=lambda: speak(greeting), daemon=True).start()
    app.mainloop()

if __name__ == "__main__":
    main()